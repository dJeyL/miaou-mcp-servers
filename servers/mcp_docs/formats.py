"""Détection de type et extraction paginée par format (PDF/docx/xlsx/pptx/zip).

`list_document` renvoie la structure sans contenu ; `read_document` renvoie un
extrait borné (`session.READ_CAP` caractères), troncature toujours signalée.
Aucun format ne renvoie jamais le document entier en un seul appel.
"""

from __future__ import annotations

import io
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path

from .search import Query, UnitResult, make_snippet, match_unit, snippets_for_unit
from .session import MAX_FILE_MB, MAX_UNZIP_MB, READ_CAP, ToolError


# ---------------------------------------------------------------------------
# Plage char/ligne relative au texte d'une unité (pdf/docx/pptx/zip texte)
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class TextRange:
    """Fenêtre de découpe dans le texte déjà produit par read() pour une unité.

    Exactement l'un des deux modes est actif (validé côté outil, pas ici) :
    caractères (`char_start` obligatoire, `char_end` optionnel) ou lignes
    (`line_start` obligatoire 1-indexé, `line_end` optionnel inclusif). La
    fenêtre déplace le point de départ ; la sortie reste plafonnée à READ_CAP
    (fenêtre glissante, pas un dump illimité, cf. décision de conception)."""

    char_start: int | None = None
    char_end: int | None = None
    line_start: int | None = None
    line_end: int | None = None

    @property
    def is_line_mode(self) -> bool:
        return self.line_start is not None

# Un membre de zip extrait en mémoire, quand il est lui-même un document
# structuré (docx/xlsx/pptx/pdf) — jamais un chemin disque.
DocSource = Path | io.BytesIO

# ---------------------------------------------------------------------------
# Détection de type — extension du nom d'origine si fourni, sinon magic bytes
# ---------------------------------------------------------------------------

_OFFICE_ZIP_SIGNATURES = {
    "word/": "docx",
    "xl/": "xlsx",
    "ppt/": "pptx",
}


def _sniff_zip_kind(source: Path | io.BytesIO) -> str:
    """Distingue docx/xlsx/pptx d'un zip brut par ses dossiers internes caractéristiques."""
    try:
        try:
            with zipfile.ZipFile(source) as z:
                names = z.namelist()
        except zipfile.BadZipFile:
            return "zip"
        for prefix, kind in _OFFICE_ZIP_SIGNATURES.items():
            if any(n.startswith(prefix) for n in names):
                return kind
        return "zip"
    finally:
        if isinstance(source, io.BytesIO):
            source.seek(0)


def detect_kind(path: Path, filename: str | None = None) -> str:
    """Renvoie l'un de : pdf, docx, xlsx, pptx, zip, inconnu.

    Le contrat client (dispatcher MIAOU) ne transmet pas le nom de fichier
    d'origine à ce jour — `filename` reste optionnel pour un usage futur ou
    manuel. En son absence, la détection retombe sur les magic bytes."""
    if filename:
        ext = Path(filename).suffix.lower().lstrip(".")
        if ext in ("pdf", "docx", "xlsx", "pptx", "zip"):
            return ext

    with open(path, "rb") as f:
        head = f.read(8)

    if head.startswith(b"%PDF"):
        return "pdf"
    if head.startswith(b"PK\x03\x04"):
        return _sniff_zip_kind(path)
    return "inconnu"


def detect_kind_from_bytes(data: bytes, filename: str | None = None) -> str:
    """Équivalent de `detect_kind` pour un contenu déjà en mémoire (membre de
    zip extrait) — même logique magic bytes / sniff zip interne, sans fichier."""
    if filename:
        ext = Path(filename).suffix.lower().lstrip(".")
        if ext in ("pdf", "docx", "xlsx", "pptx", "zip"):
            return ext

    head = data[:8]
    if head.startswith(b"%PDF"):
        return "pdf"
    if head.startswith(b"PK\x03\x04"):
        return _sniff_zip_kind(io.BytesIO(data))
    return "inconnu"


# ---------------------------------------------------------------------------
# Cap de troncature partagé par tous les formats
# ---------------------------------------------------------------------------

def _cap_text(text: str, cap: int) -> tuple[str, bool]:
    if len(text) <= cap:
        return text, False
    return text[:cap], True


def _truncation_notice(continue_hint: str) -> str:
    return f"\n\n[Tronqué à {READ_CAP} caractères — {continue_hint}]"


def _pagination_notice(continue_hint: str, unit: str = "Page") -> str:
    """Notice pour une pagination par unité (page/slide), pas par caractère :
    rien n'a été tronqué au cap, seule une unité a été servie (F6)."""
    return f"\n\n[{unit} 1 seule servie — {continue_hint}]"


def _apply_range(body: str, rng: TextRange) -> tuple[str, str]:
    """Découpe `body` selon la fenêtre char/ligne, plafonne à READ_CAP, et
    renvoie (texte, notice). La notice indique le prochain offset à demander
    quand il reste du texte au-delà de la fenêtre servie.

    Contrat de validation (assuré par l'appelant) : un seul mode actif,
    `char_start`/`line_start` non nul dans son mode."""
    if rng.is_line_mode:
        lines = body.split("\n")
        total = len(lines)
        start = rng.line_start  # 1-indexé, garanti >= 1 par la validation
        if start > total:
            return "", f"\n\n[Ligne {start} au-delà de la fin ({total} ligne(s))]"
        end = rng.line_end if rng.line_end is not None else total
        end = min(end, total)
        if end < start:
            raise ToolError(f"Plage de lignes invalide : {start}-{rng.line_end}")
        window = "\n".join(lines[start - 1 : end])
        capped, truncated = _cap_text(window, READ_CAP)
        if truncated:
            # Le cap a coupé au milieu de la fenêtre de lignes demandée. Le
            # nombre de lignes entières effectivement servies borne le prochain
            # départ (la ligne partielle finale sera relue en entier).
            served_lines = capped.count("\n")
            next_line = start + served_lines
            if served_lines == 0:
                # Le cap coupe avant la première fin de ligne : suggérer le
                # même line_start n'avancerait pas (F7) — le mode caractère
                # est la seule façon de progresser sur une ligne unique
                # dépassant le cap.
                notice = _truncation_notice(
                    f"cette ligne dépasse {READ_CAP} caractères — utiliser char_start "
                    f"plutôt que line_start pour la paginer"
                )
                return capped, notice
            notice = _truncation_notice(
                f"line_start={next_line} pour la suite (réduire la fenêtre si une "
                f"ligne unique dépasse {READ_CAP} caractères)"
            )
        elif end < total:
            notice = (
                f"\n\n[Fenêtre lignes {start}-{end} sur {total} — "
                f"line_start={end + 1} pour la suite]"
            )
        else:
            notice = ""
        return capped, notice

    # Mode caractères.
    total = len(body)
    start = rng.char_start  # garanti >= 0 par la validation
    if start >= total and total > 0:
        return "", f"\n\n[Offset {start} au-delà de la fin ({total} caractères)]"
    end = rng.char_end if rng.char_end is not None else total
    end = min(end, total)
    if end < start:
        raise ToolError(f"Plage de caractères invalide : {start}-{rng.char_end}")
    window = body[start:end]
    capped, truncated = _cap_text(window, READ_CAP)
    served_end = start + len(capped)
    if truncated or served_end < total:
        notice = _truncation_notice(f"char_start={served_end} pour la suite")
    else:
        notice = ""
    return capped, notice


def _parse_range(selector: str, total: int) -> tuple[int, int]:
    """Parse 'N' ou 'N-M' (1-indexé, inclusif), borné à [1, total]."""
    try:
        if "-" in selector:
            start_s, end_s = selector.split("-", 1)
            start, end = int(start_s), int(end_s)
        else:
            start = end = int(selector)
    except ValueError:
        raise ToolError(f"selector invalide : '{selector}' (attendu 'N' ou 'N-M')")
    start = max(1, start)
    end = min(total, end)
    if start > end:
        raise ToolError(f"selector invalide : '{selector}' (document de {total} unité(s))")
    return start, end


# ---------------------------------------------------------------------------
# PDF (pymupdf)
# ---------------------------------------------------------------------------

def _fitz_open(source: DocSource):
    import fitz

    if isinstance(source, io.BytesIO):
        return fitz.open(stream=source.getvalue(), filetype="pdf")
    return fitz.open(source)


def pdf_list(path: DocSource) -> str:
    with _fitz_open(path) as doc:
        toc = doc.get_toc()
        lines = [f"PDF — {doc.page_count} page(s)"]
        if toc:
            lines.append("Sommaire :")
            for level, title, page in toc:
                lines.append(f"{'  ' * (level - 1)}- p.{page} {title}")
        else:
            lines.append("(pas de sommaire)")
        return "\n".join(lines)


def pdf_read(path: DocSource, selector: str | None, rng: TextRange | None = None) -> str:
    with _fitz_open(path) as doc:
        if doc.page_count == 0:
            return "[PDF sans aucune page]"
        if selector:
            start, end = _parse_range(selector, doc.page_count)
        else:
            start, end = 1, 1

        parts = []
        empty_pages = []
        for page_num in range(start, end + 1):
            text = doc[page_num - 1].get_text().strip()
            if not text:
                empty_pages.append(page_num)
            parts.append(f"--- Page {page_num} ---\n{text}")
        body = "\n\n".join(parts)

        empty_note = ""
        if empty_pages:
            empty_note = (
                f"\n\n[Pages sans texte extractible (probablement scannées, "
                f"pas d'OCR en v1) : {', '.join(map(str, empty_pages))}]"
            )

        if rng is not None:
            capped, notice = _apply_range(body, rng)
            return capped + empty_note + notice

        notice = empty_note
        if not selector and doc.page_count > 1:
            notice += _pagination_notice(f"selector='2-{doc.page_count}' pour la suite")

        capped, truncated = _cap_text(body, READ_CAP)
        if truncated and not notice:
            notice = _truncation_notice("relire avec un selector plus étroit")
        return capped + notice


def pdf_search(path: DocSource, query: Query) -> list[UnitResult]:
    """Unité = page (label = numéro, round-trippable comme selector `read`)."""
    results = []
    with _fitz_open(path) as doc:
        for page_num in range(1, doc.page_count + 1):
            text = doc[page_num - 1].get_text()
            snippets = snippets_for_unit(text, query)
            if snippets:
                results.append(UnitResult(label=str(page_num), hit_count=len(snippets), snippets=snippets))
    return results


# ---------------------------------------------------------------------------
# XLSX (openpyxl)
# ---------------------------------------------------------------------------

MAX_XLSX_ROWS_DEFAULT = 200


def _open_binary(source: DocSource):
    """Fichier ouvert et passé en objet binaire, pas en chemin : openpyxl valide
    l'extension du *chemin* (rejette tout sauf .xlsx/.xlsm/...), alors que le
    fichier matérialisé par mcp_docs garde un nom stable `att-N.bin` (type
    réel détecté par magic bytes, cf. session.ref_path). Un membre de zip
    imbriqué arrive déjà en BytesIO, réutilisé tel quel."""
    if isinstance(source, io.BytesIO):
        source.seek(0)
        return source
    return open(source, "rb")


def xlsx_list(path: DocSource) -> str:
    import openpyxl

    is_owned_file = not isinstance(path, io.BytesIO)
    f = _open_binary(path)
    try:
        wb = openpyxl.load_workbook(f, read_only=True, data_only=True)
        try:
            lines = ["XLSX — feuilles :"]
            for name in wb.sheetnames:
                ws = wb[name]
                dim = ws.calculate_dimension()
                lines.append(f"- {name} ({dim}, {ws.max_row} ligne(s) x {ws.max_column} colonne(s))")
            return "\n".join(lines)
        finally:
            wb.close()
    finally:
        if is_owned_file:
            f.close()


def xlsx_read(path: DocSource, selector: str | None) -> str:
    """selector : 'NomFeuille' ou 'NomFeuille!A1:C10'. Sans plage, cap à
    MAX_XLSX_ROWS_DEFAULT lignes depuis le début de la feuille."""
    import openpyxl

    is_owned_file = not isinstance(path, io.BytesIO)
    f = _open_binary(path)
    try:
        wb = openpyxl.load_workbook(f, read_only=True, data_only=True)
        try:
            if not selector:
                sheet_name = wb.sheetnames[0]
                cell_range = None
            elif "!" in selector:
                sheet_name, cell_range = selector.split("!", 1)
            else:
                sheet_name, cell_range = selector, None

            if sheet_name not in wb.sheetnames:
                raise ToolError(f"Feuille inconnue : '{sheet_name}' (disponibles : {wb.sheetnames})")
            ws = wb[sheet_name]

            truncated_rows = False
            if cell_range:
                try:
                    selected = ws[cell_range]
                except (ValueError, IndexError) as e:
                    raise ToolError(f"Plage de cellules invalide : '{cell_range}' ({e})")
                # ws["B1"] renvoie une Cell unique (pas une plage) ; ws["A1:C10"]
                # renvoie un tuple de tuples de lignes. Normaliser au 2e cas pour
                # que le rendu ligne par ligne ci-dessous fonctionne dans les deux.
                rows = [(selected,)] if not isinstance(selected, tuple) else list(selected)
            else:
                rows = []
                for i, row in enumerate(ws.iter_rows(), start=1):
                    if i > MAX_XLSX_ROWS_DEFAULT:
                        truncated_rows = True
                        break
                    rows.append(row)

            lines = [f"Feuille '{sheet_name}' :"]
            for row in rows:
                values = [str(c.value) if c.value is not None else "" for c in row]
                lines.append(" | ".join(values))
            body = "\n".join(lines)

            notice = ""
            if truncated_rows:
                notice = (
                    f"\n\n[Tronqué à {MAX_XLSX_ROWS_DEFAULT} lignes — préciser un "
                    f"selector '{sheet_name}!A1:...' pour la suite]"
                )
            capped, capped_flag = _cap_text(body, READ_CAP)
            if capped_flag and not notice:
                notice = _truncation_notice("relire avec un selector plus étroit")
            return capped + notice
        finally:
            wb.close()
    finally:
        if is_owned_file:
            f.close()


def xlsx_search(path: DocSource, query: Query) -> list[UnitResult]:
    """Unité de résultat = feuille, mais chaque hit est une **cellule** dont le
    label (`Feuille!Coord`) est un selector `read` exact — granularité plus
    utile au modèle qu'une ligne entière. Balaie **toutes** les lignes, sans
    hériter de MAX_XLSX_ROWS_DEFAULT (sinon un match après la ligne 200 serait
    invisible, cf. AUDIT §3). Chaque cellule est matchée individuellement :
    une requête multi-termes doit avoir tous ses termes dans la MÊME cellule
    (pas de match cross-cellules)."""
    import openpyxl

    is_owned_file = not isinstance(path, io.BytesIO)
    f = _open_binary(path)
    try:
        wb = openpyxl.load_workbook(f, read_only=True, data_only=True)
        try:
            results = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                snippets = []
                occurrences = 0
                for row in ws.iter_rows():
                    for cell in row:
                        if cell.value is None:
                            continue
                        text = str(cell.value)
                        hits = match_unit(text, query)
                        if hits:
                            # Un snippet par cellule matchée (le label `Feuille!Coord`
                            # reste un selector `read` atomique), mais toutes les
                            # occurrences de la cellule comptent dans hit_count.
                            occurrences += len(hits)
                            snippet = make_snippet(text, hits[0].offset, hits[0].length)
                            snippets.append(f"{sheet_name}!{cell.coordinate} : {snippet}")
                if snippets:
                    results.append(UnitResult(label=sheet_name, hit_count=occurrences, snippets=snippets))
            return results
        finally:
            wb.close()
    finally:
        if is_owned_file:
            f.close()


# ---------------------------------------------------------------------------
# DOCX (python-docx)
# ---------------------------------------------------------------------------

_HEADING_RE = re.compile(r"^(?:Heading|Titre) (\d+)$")


def _table_text(table) -> str:
    """Concatène le texte de toutes les cellules d'une table, ligne par ligne."""
    lines = []
    for row in table.rows:
        cells = [c.text.strip() for c in row.cells]
        lines.append(" | ".join(cells))
    return "\n".join(lines)


def docx_list(path: DocSource) -> str:
    import docx

    d = docx.Document(path if isinstance(path, io.BytesIO) else str(path))
    lines = ["DOCX — sections :"]
    found = False
    for para in d.paragraphs:
        m = _HEADING_RE.match(para.style.name or "" if para.style else "")
        if m and para.text.strip():
            level = int(m.group(1))
            lines.append(f"{'  ' * (level - 1)}- {para.text.strip()}")
            found = True
    if not found:
        lines.append("(pas de heading — document sans structure de titres)")
    if d.tables:
        lines.append(f"\nTableaux ({len(d.tables)}) :")
        for i, table in enumerate(d.tables):
            lines.append(f"- table {i + 1} : {len(table.rows)} ligne(s) x {len(table.columns)} colonne(s)")
    return "\n".join(lines)


def docx_read(path: DocSource, selector: str | None, rng: TextRange | None = None) -> str:
    """selector : titre exact d'un heading (lit jusqu'au prochain heading de
    même niveau ou supérieur, paragraphes uniquement) ; sans selector, les N
    premiers paragraphes PUIS le texte de toutes les tables (un document
    purement tabulaire — formulaire questions/réponses en table, par exemple —
    n'a aucun paragraphe non-vide : sans les tables ici, `read` renverrait une
    chaîne vide alors que `list` a bien signalé leur présence)."""
    import docx

    d = docx.Document(path if isinstance(path, io.BytesIO) else str(path))
    paragraphs = d.paragraphs

    if selector:
        start_idx = None
        start_level = None
        for i, para in enumerate(paragraphs):
            m = _HEADING_RE.match(para.style.name or "" if para.style else "")
            if m and para.text.strip() == selector:
                start_idx = i
                start_level = int(m.group(1))
                break
        if start_idx is None:
            raise ToolError(f"Heading introuvable : '{selector}'")

        end_idx = len(paragraphs)
        for i in range(start_idx + 1, len(paragraphs)):
            style = paragraphs[i].style
            m = _HEADING_RE.match(style.name or "" if style else "")
            if m and int(m.group(1)) <= start_level:
                end_idx = i
                break
        selected = paragraphs[start_idx:end_idx]
        body = "\n".join(p.text for p in selected)
        notice = ""
    else:
        selected = paragraphs[:50]
        para_notice = _truncation_notice("selector=<titre d'un heading>") if len(paragraphs) > 50 else ""
        parts = [p.text for p in selected]
        if d.tables:
            parts.append("\n--- Tableaux ---")
            parts.extend(_table_text(t) for t in d.tables)
        body = "\n".join(parts)
        notice = para_notice

    if rng is not None:
        capped, range_notice = _apply_range(body, rng)
        return capped + range_notice

    capped, truncated = _cap_text(body, READ_CAP)
    if truncated and not notice:
        notice = _truncation_notice("relire avec un selector plus étroit")
    return capped + notice


_PREAMBLE_LABEL = "(préambule)"
_BODY_LABEL = "(corps)"
_TABLES_LABEL = "(tableaux)"


def _docx_sections(d) -> list[tuple[str, str]]:
    """Découpe un document docx en unités `(label, texte)` réutilisables pour
    `docx_search`. Une section = un heading et tout le texte jusqu'au prochain
    heading de même niveau ou supérieur (même règle de bornage que
    `docx_read`, dupliquée ici plutôt que factorisée pour ne pas risquer de
    régression sur `docx_read`, déjà couvert par ses propres tests).

    Le texte avant le premier heading (ou tout le document s'il n'y a aucun
    heading) n'a pas de label round-trippable naturel : il retombe sur
    `(préambule)` s'il précède un heading existant, ou `(corps)` si le
    document n'a aucun heading — `read` sans selector renvoie ce même
    contenu dans les deux cas, donc le round-trip reste valide (round-trip
    partiel assumé, cf. PLAN décision #2). Les tables n'appartiennent à aucun
    heading : regroupées sous `(tableaux)`, round-trippable de la même façon
    (read sans selector les inclut aussi)."""
    paragraphs = d.paragraphs
    headings = []  # (idx, level, title)
    for i, para in enumerate(paragraphs):
        m = _HEADING_RE.match(para.style.name or "" if para.style else "")
        if m and para.text.strip():
            headings.append((i, int(m.group(1)), para.text.strip()))

    sections: list[tuple[str, str]] = []

    if not headings:
        body = "\n".join(p.text for p in paragraphs)
        sections.append((_BODY_LABEL, body))
    else:
        first_idx = headings[0][0]
        if first_idx > 0:
            preamble = "\n".join(p.text for p in paragraphs[:first_idx])
            if preamble.strip():
                sections.append((_PREAMBLE_LABEL, preamble))

        for pos, (start_idx, level, title) in enumerate(headings):
            end_idx = len(paragraphs)
            for next_idx, next_level, _ in headings[pos + 1 :]:
                if next_level <= level:
                    end_idx = next_idx
                    break
            text = "\n".join(p.text for p in paragraphs[start_idx:end_idx])
            sections.append((title, text))

    if d.tables:
        sections.append((_TABLES_LABEL, "\n".join(_table_text(t) for t in d.tables)))

    return sections


def docx_search(path: DocSource, query: Query) -> list[UnitResult]:
    """Unité = section heading (label = titre exact, selector `read` valide) ;
    texte hors-section et docs sans heading → labels spéciaux `(préambule)`/
    `(corps)`, tables → `(tableaux)`. Round-trip partiel assumé pour ces
    labels spéciaux (décision #2 : `read` sans selector les re-sert)."""
    import docx

    d = docx.Document(path if isinstance(path, io.BytesIO) else str(path))
    results = []
    for label, text in _docx_sections(d):
        snippets = snippets_for_unit(text, query)
        if snippets:
            results.append(UnitResult(label=label, hit_count=len(snippets), snippets=snippets))
    return results


# ---------------------------------------------------------------------------
# PPTX (python-pptx)
# ---------------------------------------------------------------------------

def pptx_list(path: DocSource) -> str:
    from pptx import Presentation

    prs = Presentation(path if isinstance(path, io.BytesIO) else str(path))
    lines = [f"PPTX — {len(prs.slides)} slide(s) :"]
    for i, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text if slide.shapes.title else "(sans titre)"
        lines.append(f"{i}. {title}")
    return "\n".join(lines)


def pptx_read(path: DocSource, selector: str | None, rng: TextRange | None = None) -> str:
    from pptx import Presentation

    prs = Presentation(path if isinstance(path, io.BytesIO) else str(path))
    total = len(prs.slides)
    if total == 0:
        return "[PPTX sans aucune slide]"

    if selector:
        start, end = _parse_range(selector, total)
    else:
        start, end = 1, 1

    parts = []
    for slide_num in range(start, end + 1):
        slide = prs.slides[slide_num - 1]
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(t for t in texts if t.strip()))
    body = "\n\n".join(parts)

    if rng is not None:
        capped, notice = _apply_range(body, rng)
        return capped + notice

    notice = ""
    if not selector and total > 1:
        notice = _pagination_notice(f"selector='2-{total}' pour la suite", unit="Slide")
    capped, truncated = _cap_text(body, READ_CAP)
    if truncated and not notice:
        notice = _truncation_notice("relire avec un selector plus étroit")
    return capped + notice


def pptx_search(path: DocSource, query: Query) -> list[UnitResult]:
    """Unité = slide (label = numéro, round-trippable comme selector `read`)."""
    from pptx import Presentation

    prs = Presentation(path if isinstance(path, io.BytesIO) else str(path))
    results = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
        text = "\n".join(t for t in texts if t.strip())
        snippets = snippets_for_unit(text, query)
        if snippets:
            results.append(UnitResult(label=str(slide_num), hit_count=len(snippets), snippets=snippets))
    return results


# ---------------------------------------------------------------------------
# ZIP (stdlib zipfile) — arbre d'entrées et lecture de membre
#
# Gardes de sécurité : zip-slip (chemin résolu, rejet absolu/`..`), tailles
# totale/par-entrée contrôlées en flux (les headers mentent — on ne se fie pas
# à ZipInfo.file_size seul), rejet des membres chiffrés. Un membre qui est
# lui-même un document structuré (pdf/docx/xlsx/pptx/zip) est extractible via
# dispatch récursif borné à un seul niveau d'imbrication (read_nested_member /
# list_nested_member). Menace : utilisateur unique, mais les gardes restent
# non négociables (coût trivial, échec = dommage filesystem).
# ---------------------------------------------------------------------------

_ZIP_READ_CHUNK = 64 * 1024


def _is_zip_slip(member_path: str) -> bool:
    """Rejette chemin absolu ou toute composante '..' (avant même resolve())."""
    if member_path.startswith("/") or member_path.startswith("\\"):
        return True
    normalized = Path(member_path)
    return ".." in normalized.parts


def _is_encrypted(info: zipfile.ZipInfo) -> bool:
    return bool(info.flag_bits & 0x1)


def _is_nested_archive(head: bytes) -> bool:
    return head.startswith(b"PK\x03\x04")


def zip_list(path: DocSource) -> str:
    try:
        with zipfile.ZipFile(path) as z:
            infos = z.infolist()
    except zipfile.BadZipFile as e:
        raise ToolError(f"Archive zip corrompue ou invalide : {e}")

    lines = ["ZIP — entrées :"]
    total_declared = 0
    for info in infos:
        kind = "dir" if info.is_dir() else "file"
        flags = []
        if _is_zip_slip(info.filename):
            flags.append("chemin suspect, non extractible")
        if not info.is_dir() and _is_encrypted(info):
            flags.append("chiffré, non extractible")
        if not info.is_dir() and info.filename.lower().endswith((".docx", ".xlsx", ".pptx", ".pdf")):
            flags.append("document structuré, lisible via read(path=...)")
        elif not info.is_dir() and info.filename.lower().endswith(".zip"):
            flags.append("archive imbriquée potentielle, non extractible au-delà d'un niveau")
        flag_note = f" [{', '.join(flags)}]" if flags else ""
        lines.append(f"- {info.filename} ({info.file_size} octets, {kind}){flag_note}")
        total_declared += info.file_size

    total_mb = total_declared / (1024 * 1024)
    if total_mb > MAX_UNZIP_MB:
        lines.append(
            f"\n[Taille décompressée totale déclarée : {total_mb:.1f} Mo, "
            f"dépasse la limite de {MAX_UNZIP_MB} Mo — l'extraction de membres "
            f"individuels reste bornée par membre, mais l'archive entière ne "
            f"pourrait pas être extraite intégralement]"
        )
    return "\n".join(lines)


def _extract_zip_member_bytes(path: DocSource, member_path: str) -> bytes:
    """Extrait un membre du zip en mémoire, en flux et avec garde de taille
    cumulative (pas de confiance dans ZipInfo.file_size). Rejette zip-slip,
    membres chiffrés et archive corrompue. Ne dit rien sur le *contenu* extrait
    (archive imbriquée ou non) — c'est aux appelants de le déterminer."""
    if _is_zip_slip(member_path):
        raise ToolError(f"Chemin de membre invalide (traversal) : '{member_path}'")

    with zipfile.ZipFile(path) as z:
        try:
            info = z.getinfo(member_path)
        except KeyError as e:
            raise ToolError(f"Membre introuvable dans l'archive : '{member_path}'") from e
        if info.is_dir():
            raise ToolError(f"'{member_path}' est un répertoire, pas un fichier")
        if _is_encrypted(info):
            raise ToolError(f"Membre chiffré, non extractible : '{member_path}'")

        # Garde rapide sur la taille déclarée (header), puis contrôle réel en
        # flux ci-dessous — les deux sont nécessaires, le header seul ment.
        if info.file_size / (1024 * 1024) > MAX_UNZIP_MB:
            raise ToolError(
                f"Membre trop volumineux d'après l'en-tête "
                f"({info.file_size / (1024 * 1024):.1f} Mo, max {MAX_UNZIP_MB} Mo)"
            )

        chunks = []
        total = 0
        try:
            with z.open(info) as member_f:
                while True:
                    chunk = member_f.read(_ZIP_READ_CHUNK)
                    if not chunk:
                        break
                    total += len(chunk)
                    if total / (1024 * 1024) > MAX_UNZIP_MB:
                        raise ToolError(
                            f"Membre trop volumineux en flux réel (>{MAX_UNZIP_MB} Mo, "
                            f"en-tête mensonger) : '{member_path}'"
                        )
                    chunks.append(chunk)
        except zipfile.BadZipFile as e:
            raise ToolError(f"Archive corrompue ou en-tête incohérent pour '{member_path}' : {e}") from e
        return b"".join(chunks)


def _check_nested_member_size(data: bytes, member_path: str) -> None:
    """Un membre imbriqué reconnu comme docx/xlsx/pptx/pdf passe la garde zip
    (MAX_UNZIP_MB, ex. 100 Mo) mais doit rester sous MAX_FILE_MB (ex. 20 Mo)
    avant d'être parsé par une lib lourde — même borne que la matérialisation
    initiale d'un attachement (session.materialize)."""
    size_mb = len(data) / (1024 * 1024)
    if size_mb > MAX_FILE_MB:
        raise ToolError(
            f"Membre imbriqué trop volumineux pour être extrait comme document "
            f"structuré ({size_mb:.1f} Mo, max {MAX_FILE_MB} Mo) : '{member_path}'"
        )


def read_nested_member(
    member_path: str, data: bytes, selector: str | None, rng: TextRange | None = None
) -> str:
    """Dispatch récursif : le membre extrait d'un zip est lui-même un document
    structuré. Un seul niveau d'imbrication est supporté (pas de zip-dans-zip-
    dans-zip) — un membre qui est lui-même une archive contenant une archive
    reste signalé mais non extrait, pour borner la récursion et le coût de
    parsing (chaque niveau individuel pourrait respecter MAX_UNZIP_MB tout en
    démultipliant le travail cumulé)."""
    kind = detect_kind_from_bytes(data, filename=member_path)
    if kind == "zip":
        return (
            f"[Archive imbriquée à plus d'un niveau détectée dans '{member_path}' "
            f"— non extractible, un seul niveau d'imbrication est supporté]"
        )
    if kind == "inconnu":
        return f"[Membre binaire non affichable en l'état : '{member_path}', {len(data)} octets]"

    _check_nested_member_size(data, member_path)
    try:
        return read_document(kind, io.BytesIO(data), selector, rng)
    except zipfile.BadZipFile as e:
        raise ToolError(f"Membre imbriqué corrompu, non extractible : '{member_path}' ({e})")


def list_nested_member(member_path: str, data: bytes) -> str:
    """Équivalent de `read_nested_member` pour `docs__list` — structure d'un
    membre imbriqué sans en renvoyer le contenu."""
    kind = detect_kind_from_bytes(data, filename=member_path)
    if kind == "zip":
        return (
            f"[Archive imbriquée à plus d'un niveau détectée dans '{member_path}' "
            f"— non listable, un seul niveau d'imbrication est supporté]"
        )
    if kind == "inconnu":
        return f"[Membre binaire non reconnu : '{member_path}', {len(data)} octets]"

    _check_nested_member_size(data, member_path)
    try:
        return list_document(kind, io.BytesIO(data))
    except zipfile.BadZipFile as e:
        raise ToolError(f"Membre imbriqué corrompu, non listable : '{member_path}' ({e})")


def _is_structured_document(head: bytes) -> bool:
    """PDF ou zip (brut/Office) — tout ce que `read_nested_member` sait dispatcher."""
    return head.startswith(b"%PDF") or _is_nested_archive(head)


def zip_read_member(
    path: DocSource,
    member_path: str,
    selector: str | None = None,
    rng: TextRange | None = None,
) -> str:
    """Rejette zip-slip et membres chiffrés (via `_extract_zip_member_bytes`).
    Si le membre est lui-même un document structuré (pdf/docx/xlsx/pptx/zip),
    dispatch récursif borné à un niveau (`read_nested_member`, `selector` et
    `rng` transmis) ; sinon traité comme texte brut (`selector` alors sans
    effet, `rng` applique la fenêtre char/ligne sur le texte du membre)."""
    data = _extract_zip_member_bytes(path, member_path)

    if _is_structured_document(data[:8]):
        return read_nested_member(member_path, data, selector, rng)

    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        return f"[Membre binaire non affichable en l'état : '{member_path}', {len(data)} octets]"

    if rng is not None:
        capped, notice = _apply_range(text, rng)
        return capped + notice

    capped, truncated = _cap_text(text, READ_CAP)
    notice = _truncation_notice("relire avec un selector plus étroit") if truncated else ""
    return capped + notice


def zip_extract_member_text(path: DocSource, member_path: str) -> str:
    """Extrait le **texte intégral** d'un membre zip, sans `_apply_range`/`READ_CAP`
    — cette voie ne renvoie rien au contexte du modèle : les octets transitent vers
    le client (canal `content_b64`/`res_…`, lot K MIAOU) qui les matérialise en
    ressource adressable par `js__eval`. `READ_CAP` borne le *contexte*, pas ce
    *transfert*. Mêmes gardes que `zip_read_member` (zip-slip, chiffré, taille en
    flux, via `_extract_zip_member_bytes`) ; un membre structuré (pdf/docx/xlsx/
    pptx/zip imbriqué) est refusé explicitement — `docs__extract` ne cible que le
    texte brut (logs/JSON/CSV), `docs__read` reste la voie pour un membre structuré."""
    data = _extract_zip_member_bytes(path, member_path)

    if _is_structured_document(data[:8]):
        raise ToolError(
            f"'{member_path}' est un document structuré (pdf/docx/xlsx/pptx/zip), "
            f"pas un membre texte — docs__extract ne cible que le texte brut. "
            f"Utiliser docs__read(path='{member_path}') ou docs__list(path='{member_path}')."
        )

    try:
        return data.decode("utf-8")
    except UnicodeDecodeError as e:
        raise ToolError(f"Membre binaire non décodable en UTF-8 : '{member_path}'") from e


def zip_search(path: DocSource, query: Query, member_path: str | None = None) -> list[UnitResult]:
    """Cherche dans les membres texte brut d'une archive zip. Sans
    `member_path`, balaie tous les membres non ignorés ; avec `member_path`,
    restreint la recherche à ce seul membre (décision #6). Un membre
    reconnu comme document structuré (docx/xlsx/pptx/pdf), chiffré, une
    archive imbriquée ou un fichier binaire non-utf8 est ignoré et listé en
    note finale (décision #3) — pas de dispatch récursif ici, contrairement à
    `zip_read_member` (limitation explicite du brief, zip = texte seul pour
    la recherche). Label unité = chemin du membre (round-trip :
    `read(path=member_path)`)."""
    if member_path is not None:
        # Ciblage explicite : erreur claire si le membre n'existe pas ou est
        # inextractible, comme zip_read_member (pas un skip silencieux ici).
        data = _extract_zip_member_bytes(path, member_path)
        if _is_structured_document(data[:8]):
            note = (
                f"[Membre '{member_path}' non cherché : document structuré ou "
                f"binaire, la recherche zip ne porte que sur le texte brut]"
            )
            return [UnitResult(label="(note)", hit_count=0, snippets=[note])]
        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            note = f"[Membre '{member_path}' non cherché : contenu binaire]"
            return [UnitResult(label="(note)", hit_count=0, snippets=[note])]
        snippets = snippets_for_unit(text, query)
        if not snippets:
            return []
        return [UnitResult(label=member_path, hit_count=len(snippets), snippets=snippets)]

    results = []
    skipped: list[str] = []
    try:
        with zipfile.ZipFile(path) as z:
            infos = z.infolist()
    except zipfile.BadZipFile as e:
        raise ToolError(f"Archive zip corrompue ou invalide : {e}")

    for info in infos:
        if info.is_dir():
            continue
        name = info.filename
        if _is_zip_slip(name):
            skipped.append(f"{name} (chemin suspect)")
            continue
        if _is_encrypted(info):
            skipped.append(f"{name} (chiffré)")
            continue
        if info.file_size / (1024 * 1024) > MAX_UNZIP_MB:
            skipped.append(f"{name} (trop volumineux)")
            continue

        try:
            data = _extract_zip_member_bytes(path, name)
        except ToolError:
            skipped.append(f"{name} (illisible)")
            continue

        if _is_structured_document(data[:8]):
            skipped.append(f"{name} (document structuré, non cherché)")
            continue

        try:
            text = data.decode("utf-8")
        except UnicodeDecodeError:
            skipped.append(f"{name} (binaire)")
            continue

        snippets = snippets_for_unit(text, query)
        if snippets:
            results.append(UnitResult(label=name, hit_count=len(snippets), snippets=snippets))

    if skipped:
        note = "[Membres non cherchés (zones aveugles) : " + ", ".join(skipped) + "]"
        results.append(UnitResult(label="(note)", hit_count=0, snippets=[note]))

    return results


def zip_list_member(path: DocSource, member_path: str) -> str:
    """Liste la structure d'un membre imbriqué (docx/xlsx/pptx/pdf) sans
    matérialisation disque. Un membre non reconnu comme document structuré (ou
    une archive imbriquée à plus d'un niveau) reste signalé, jamais planté."""
    data = _extract_zip_member_bytes(path, member_path)

    if _is_structured_document(data[:8]):
        return list_nested_member(member_path, data)

    return f"[Membre '{member_path}' : {len(data)} octets, pas un document structuré (docx/xlsx/pptx/pdf/zip)]"


# ---------------------------------------------------------------------------
# Dispatch — table kind → fonction
# ---------------------------------------------------------------------------

_LIST_DISPATCH = {
    "pdf": pdf_list,
    "docx": docx_list,
    "xlsx": xlsx_list,
    "pptx": pptx_list,
    "zip": zip_list,
}

_READ_DISPATCH = {
    "pdf": pdf_read,
    "docx": docx_read,
    "xlsx": xlsx_read,
    "pptx": pptx_read,
}

_SEARCH_DISPATCH = {
    "pdf": pdf_search,
    "xlsx": xlsx_search,
    "docx": docx_search,
    "pptx": pptx_search,
}


def list_document(kind: str, path: DocSource) -> str:
    fn = _LIST_DISPATCH.get(kind)
    if fn is None:
        raise ToolError(f"Type de document non supporté ou non reconnu (détecté : '{kind}')")
    return fn(path)


def read_document(kind: str, path: DocSource, selector: str | None, rng: TextRange | None = None) -> str:
    if kind == "zip":
        raise ToolError("Zip : préciser 'path' pour lire un membre (voir docs__list pour l'arbre)")
    if rng is not None and kind == "xlsx":
        raise ToolError(
            "Plage char/ligne non applicable à un xlsx (grille, pas de flux texte) — "
            "utiliser un selector 'Feuille!A1:C10'"
        )
    fn = _READ_DISPATCH.get(kind)
    if fn is None:
        raise ToolError(f"Type de document non supporté ou non reconnu (détecté : '{kind}')")
    # xlsx_read n'accepte pas rng (rejeté ci-dessus) : appel à 2 arguments.
    if kind == "xlsx":
        return fn(path, selector)
    return fn(path, selector, rng)


def search_document(kind: str, path: DocSource, query: Query) -> list[UnitResult]:
    fn = _SEARCH_DISPATCH.get(kind)
    if fn is None:
        raise ToolError(f"Recherche non encore supportée pour ce format (détecté : '{kind}')")
    return fn(path, query)
