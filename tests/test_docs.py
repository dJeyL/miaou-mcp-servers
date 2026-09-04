"""Tests pour servers/mcp_docs/ (D1 — sessions/transport/REF_UNKNOWN, D2 — formats,
D3 — sécurité archives)."""
import base64
import shutil
import sys
import time
from pathlib import Path
from unittest.mock import patch

import pytest

_HAS_ZIP_CLI = shutil.which("zip") is not None
_requires_zip_cli = pytest.mark.skipif(
    not _HAS_ZIP_CLI, reason="binaire 'zip' absent (stdlib zipfile ne sait pas chiffrer en écriture)"
)

_ROOT = Path(__file__).parent.parent
_SERVERS = _ROOT / "servers"
for p in (_ROOT, _SERVERS):
    if str(p) not in sys.path:
        sys.path.insert(0, str(p))

import mcp_docs
from mcp_docs import REF_UNKNOWN_ERROR_CODE, REF_UNKNOWN_SENTINEL, ToolError
from mcp_docs import session as docs_session
from mcp_docs.formats import (
    detect_kind,
    list_document,
    read_document,
    zip_extract_member_text,
    zip_list_member,
    zip_read_member,
)
from mcp_docs.search import Query, fold, make_snippet, match_unit, parse_query
from mcp_docs.session import (
    materialize,
    resolve_ref,
    session_dir,
    sweep_expired_sessions,
    validate_ref,
    validate_session_id,
)


@pytest.fixture()
def workdir(tmp_path, monkeypatch):
    monkeypatch.setattr(docs_session, "WORKDIR", tmp_path)
    return tmp_path


# ---------------------------------------------------------------------------
# Sanitization session_id / ref
# ---------------------------------------------------------------------------

def test_validate_session_id_rejects_empty():
    with pytest.raises(ToolError):
        validate_session_id("")


def test_validate_session_id_rejects_path_separator():
    with pytest.raises(ToolError):
        validate_session_id("foo/bar")
    with pytest.raises(ToolError):
        validate_session_id("foo\\bar")


def test_validate_session_id_rejects_dotdot():
    with pytest.raises(ToolError):
        validate_session_id("../etc")
    with pytest.raises(ToolError):
        validate_session_id("foo..bar")


def test_validate_session_id_accepts_normal_id():
    assert validate_session_id("conv-123") == "conv-123"


def test_validate_session_id_rejects_dot(workdir):
    """D1 : "." == WORKDIR lui-même — sans ce rejet, drop_session(".") efface tout."""
    with pytest.raises(ToolError):
        validate_session_id(".")


def test_validate_session_id_rejects_nul_byte():
    """DOC8 : un NUL non filtré passe validate_session_id puis fait fuiter
    ValueError('embedded null character in path') à mkdir dans touch_session."""
    with pytest.raises(ToolError):
        validate_session_id("a\x00b")


def test_docs_env_int_invalid_value_raises_named_error(monkeypatch):
    """DOC5 : même comportement que mcp_web.cache._env_int — nommer la
    variable plutôt qu'un ValueError anonyme à l'import."""
    monkeypatch.setenv("MIAOU_DOCS_READ_CAP", "not-a-number")
    with pytest.raises(ValueError, match="MIAOU_DOCS_READ_CAP"):
        docs_session._env_int("MIAOU_DOCS_READ_CAP", 20000)


def test_validate_ref_accepts_att_n():
    assert validate_ref("att-1") == "att-1"
    assert validate_ref("att-42") == "att-42"


def test_validate_ref_accepts_file_id():
    """file-<id> : fichier de bibliothèque d'espace (MIAOU lot Cbis,
    libraryRefFromId) — même statut que att-N, pas une forme dégradée."""
    assert validate_ref("file-2rhku6t4") == "file-2rhku6t4"
    assert validate_ref("file-a1b2c3") == "file-a1b2c3"


def test_validate_ref_rejects_file_id_uppercase_or_empty_suffix():
    with pytest.raises(ToolError):
        validate_ref("file-ABC123")
    with pytest.raises(ToolError):
        validate_ref("file-")


def test_validate_ref_accepts_res_id():
    """res_<id> : ressource de session côté client (MIAOU lot K,
    generateResourceId) — underscore après "res", pas un tiret comme
    att-/file-."""
    assert validate_ref("res_abc123") == "res_abc123"
    assert validate_ref("res_2rhku6t4") == "res_2rhku6t4"


def test_validate_ref_rejects_res_id_with_dash_or_uppercase():
    with pytest.raises(ToolError):
        validate_ref("res-abc")
    with pytest.raises(ToolError):
        validate_ref("res_ABC123")
    with pytest.raises(ToolError):
        validate_ref("res_")


def test_validate_ref_rejects_nested_path():
    """La syntaxe ref#path du brief est explicitement écartée (audit §2) —
    l'adressage de membre d'archive passe par le paramètre `path` séparé."""
    with pytest.raises(ToolError):
        validate_ref("att-3#membre.docx")


def test_validate_ref_rejects_garbage():
    with pytest.raises(ToolError):
        validate_ref("not-a-ref")


# ---------------------------------------------------------------------------
# Matérialisation idempotente
# ---------------------------------------------------------------------------

def test_materialize_writes_file(workdir):
    content = base64.b64encode(b"hello world").decode()
    path = materialize("conv-1", "att-1", content)
    assert path.exists()
    assert path.read_bytes() == b"hello world"


def test_materialize_is_idempotent(workdir):
    """Un rechargement de page MIAOU repousse content_b64 pour un ref déjà connu —
    ré-acceptation silencieuse SANS réécriture (le fichier déjà matérialisé
    fait foi), pas une erreur (T6 : la docstring recopiait auparavant la
    formulation « réécriture silencieuse », inexacte — corrigée avec DD2)."""
    content = base64.b64encode(b"hello world").decode()
    path1 = materialize("conv-1", "att-1", content)
    path1.write_bytes(b"already there, must not be overwritten")

    other_content = base64.b64encode(b"different bytes").decode()
    path2 = materialize("conv-1", "att-1", other_content)

    assert path1 == path2
    assert path2.read_bytes() == b"already there, must not be overwritten"


def test_materialize_rejects_oversized_file(workdir):
    with patch.object(docs_session, "MAX_FILE_MB", 1):
        big = base64.b64encode(b"x" * (2 * 1024 * 1024)).decode()
        with pytest.raises(ToolError, match="volumineux"):
            materialize("conv-1", "att-1", big)


def test_materialize_rejects_session_quota(workdir):
    with patch.object(docs_session, "MAX_SESSION_MB", 1), patch.object(docs_session, "MAX_FILE_MB", 10):
        # Premier fichier ~0.9 Mo passe.
        first = base64.b64encode(b"x" * int(0.9 * 1024 * 1024)).decode()
        materialize("conv-1", "att-1", first)

        # Deuxième fichier ferait dépasser le quota de session (1 Mo).
        second = base64.b64encode(b"y" * int(0.5 * 1024 * 1024)).decode()
        with pytest.raises(ToolError, match="Quota"):
            materialize("conv-1", "att-2", second)


def test_materialize_rejects_bad_session_id(workdir):
    content = base64.b64encode(b"data").decode()
    with pytest.raises(ToolError):
        materialize("../escape", "att-1", content)


def test_materialize_rejects_invalid_base64(workdir):
    """DOC2 : un base64 invalide (padding incorrect) doit lever une ToolError
    claire, pas laisser fuiter binascii.Error brute au client."""
    with pytest.raises(ToolError, match="content_b64 invalide"):
        materialize("conv-1", "att-1", "not-valid-base64!!!")


def test_materialize_writes_file_for_library_ref(workdir):
    """file-<id> (fichier de bibliothèque d'espace, MIAOU lot Cbis) suit
    exactement le même chemin que att-N — pas un cas particulier."""
    content = base64.b64encode(b"library content").decode()
    path = materialize("conv-1", "file-2rhku6t4", content)
    assert path.exists()
    assert path.read_bytes() == b"library content"


def test_materialize_writes_file_for_resource_ref(workdir):
    """res_<id> (ressource de session côté client, MIAOU lot K) suit
    exactement le même chemin que att-N/file-<id> — pas un cas particulier."""
    content = base64.b64encode(b"resource content").decode()
    path = materialize("conv-1", "res_abc123", content)
    assert path.exists()
    assert path.read_bytes() == b"resource content"


def test_materialize_att_and_file_refs_do_not_collide(workdir):
    """Un att-N et un file-<id> dans la même session matérialisent des fichiers
    distincts (préfixes différents dans ref_path) — jamais un écrasement croisé."""
    att_content = base64.b64encode(b"attachment bytes").decode()
    file_content = base64.b64encode(b"library bytes").decode()
    att_path = materialize("conv-1", "att-1", att_content)
    file_path = materialize("conv-1", "file-1", file_content)
    assert att_path != file_path
    assert att_path.read_bytes() == b"attachment bytes"
    assert file_path.read_bytes() == b"library bytes"


# ---------------------------------------------------------------------------
# resolve_ref — REF_UNKNOWN
# ---------------------------------------------------------------------------

def test_resolve_ref_unknown_without_content_raises_sentinel(workdir):
    with pytest.raises(ToolError) as exc_info:
        resolve_ref("conv-1", "att-99", None)
    assert str(exc_info.value).startswith(REF_UNKNOWN_SENTINEL)


def test_resolve_ref_materializes_when_content_provided(workdir):
    content = base64.b64encode(b"payload").decode()
    path = resolve_ref("conv-1", "att-1", content)
    assert path.read_bytes() == b"payload"


def test_resolve_ref_known_ref_without_content_succeeds(workdir):
    content = base64.b64encode(b"payload").decode()
    resolve_ref("conv-1", "att-1", content)
    path = resolve_ref("conv-1", "att-1", None)
    assert path.read_bytes() == b"payload"


# ---------------------------------------------------------------------------
# TTL sweep (horloge mockée)
# ---------------------------------------------------------------------------

def test_sweep_removes_expired_session(workdir):
    old_dir = session_dir("conv-old")
    old_dir.mkdir(parents=True)
    old_time = time.time() - 25 * 3600  # TTL défaut 24h dépassé
    import os
    os.utime(old_dir, (old_time, old_time))

    with patch.object(docs_session, "TTL_HOURS", 24):
        sweep_expired_sessions()

    assert not old_dir.exists()


def test_sweep_keeps_recent_session(workdir):
    recent_dir = session_dir("conv-recent")
    recent_dir.mkdir(parents=True)

    with patch.object(docs_session, "TTL_HOURS", 24):
        sweep_expired_sessions()

    assert recent_dir.exists()


def test_sweep_touch_on_access_prevents_expiry(workdir):
    """touch_session doit repousser mtime — un accès récent empêche le sweep."""
    d = session_dir("conv-1")
    d.mkdir(parents=True)
    old_time = time.time() - 25 * 3600
    import os
    os.utime(d, (old_time, old_time))

    docs_session.touch_session("conv-1")

    with patch.object(docs_session, "TTL_HOURS", 24):
        sweep_expired_sessions()

    assert d.exists()


# ---------------------------------------------------------------------------
# docs__drop_session
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_drop_session_removes_directory(workdir):
    from mcp_docs import server as docs_server

    d = session_dir("conv-1")
    d.mkdir(parents=True)
    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool("drop_session", {"session_id": "conv-1"})
    assert not d.exists()
    assert "conv-1" in result


@pytest.mark.asyncio
async def test_drop_session_idempotent_when_absent(workdir):
    from mcp_docs import server as docs_server

    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool("drop_session", {"session_id": "conv-never-existed"})
    assert "conv-never-existed" in result


@pytest.mark.asyncio
async def test_drop_session_sweeps_other_expired_sessions(workdir):
    """DD3/DOC7 : docs/servers.md affirme le sweep TTL 'en tête de chaque appel
    d'outil' — drop_session doit désormais le faire aussi, pas seulement
    list/read/search/extract."""
    from mcp_docs import server as docs_server

    other = session_dir("conv-expired")
    other.mkdir(parents=True)
    old_time = time.time() - 999999
    import os
    os.utime(other, (old_time, old_time))

    target = session_dir("conv-target")
    target.mkdir(parents=True)

    tm = docs_server.mcp._tool_manager
    await tm.call_tool("drop_session", {"session_id": "conv-target"})

    assert not other.exists()  # balayée par le sweep, pas par drop_session lui-même


# ---------------------------------------------------------------------------
# REF_UNKNOWN à travers le stack proxy (modèle test_proxy.py)
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_ref_unknown_through_proxy_raises_jsonrpc_error(workdir):
    """Le sentinel textuel REF_UNKNOWN doit ressortir en erreur JSON-RPC
    data.code == 'REF_UNKNOWN', pas en isError textuel (contrat client, audit §3)."""
    import mcp.types as types
    from mcp.shared.exceptions import McpError

    from mcp_proxy import InProcessUpstream, build_proxy_server

    upstream = InProcessUpstream("mcp_docs")
    await upstream.start()

    upstreams = {"docs": upstream}
    tool_map: dict = {}
    server = build_proxy_server(upstreams, tool_map)

    list_handler = server.request_handlers[types.ListToolsRequest]
    await list_handler(types.ListToolsRequest(method="tools/list", params=None))

    call_handler = server.request_handlers[types.CallToolRequest]
    request = types.CallToolRequest(
        method="tools/call",
        params=types.CallToolRequestParams(
            name="docs__list",
            arguments={"ref": "att-99", "session_id": "conv-1"},
        ),
    )

    with pytest.raises(McpError) as exc_info:
        await call_handler(request)

    assert exc_info.value.error.data["code"] == "REF_UNKNOWN"
    assert exc_info.value.error.code == REF_UNKNOWN_ERROR_CODE


@pytest.mark.asyncio
async def test_non_ref_unknown_error_stays_iserror_through_proxy(workdir):
    """Une erreur applicative ordinaire (session_id manquant) ne doit pas être
    convertie en erreur JSON-RPC — seul le sentinel REF_UNKNOWN déclenche McpError."""
    import mcp.types as types

    from mcp_proxy import InProcessUpstream, build_proxy_server

    upstream = InProcessUpstream("mcp_docs")
    await upstream.start()

    upstreams = {"docs": upstream}
    tool_map: dict = {}
    server = build_proxy_server(upstreams, tool_map)

    list_handler = server.request_handlers[types.ListToolsRequest]
    await list_handler(types.ListToolsRequest(method="tools/list", params=None))

    call_handler = server.request_handlers[types.CallToolRequest]
    request = types.CallToolRequest(
        method="tools/call",
        params=types.CallToolRequestParams(
            name="docs__list",
            arguments={"ref": "att-1"},
        ),
    )

    result = await call_handler(request)
    assert result.root.isError


# ---------------------------------------------------------------------------
# Détection de type
# ---------------------------------------------------------------------------

def test_detect_kind_by_filename_extension(tmp_path):
    p = tmp_path / "whatever.bin"
    p.write_bytes(b"not really a pdf")
    assert detect_kind(p, filename="report.pdf") == "pdf"
    assert detect_kind(p, filename="sheet.xlsx") == "xlsx"


def test_detect_kind_pdf_by_magic_bytes(tmp_path):
    p = tmp_path / "f.bin"
    p.write_bytes(b"%PDF-1.4\n...")
    assert detect_kind(p) == "pdf"


def test_detect_kind_unknown_for_garbage(tmp_path):
    p = tmp_path / "f.bin"
    p.write_bytes(b"not a document at all")
    assert detect_kind(p) == "inconnu"


def test_detect_kind_plain_zip_by_magic_bytes(tmp_path):
    import zipfile

    p = tmp_path / "f.bin"
    with zipfile.ZipFile(p, "w") as z:
        z.writestr("hello.txt", "hi")
    assert detect_kind(p) == "zip"


# ---------------------------------------------------------------------------
# PDF (pymupdf) — fixtures générées à la volée
# ---------------------------------------------------------------------------

def _make_pdf(tmp_path, pages_text, toc=None):
    import fitz

    doc = fitz.open()
    for text in pages_text:
        page = doc.new_page()
        if text:
            page.insert_text((72, 72), text)
    if toc:
        doc.set_toc(toc)
    path = tmp_path / "doc.pdf"
    doc.save(str(path))
    doc.close()
    return path


def test_pdf_list_reports_page_count_and_outline(tmp_path):
    path = _make_pdf(tmp_path, ["Page one", "Page two"], toc=[[1, "Chapter 1", 1], [1, "Chapter 2", 2]])
    result = list_document("pdf", path)
    assert "2 page" in result
    assert "Chapter 1" in result
    assert "Chapter 2" in result


def test_pdf_list_no_outline(tmp_path):
    path = _make_pdf(tmp_path, ["Solo page"])
    result = list_document("pdf", path)
    assert "pas de sommaire" in result


def test_pdf_read_without_selector_returns_first_page_only(tmp_path):
    """F6 : la notice doit refléter une pagination par page, pas une
    troncature de caractères (rien n'est tronqué ici)."""
    path = _make_pdf(tmp_path, ["First page text", "Second page text"])
    result = read_document("pdf", path, None)
    assert "First page text" in result
    assert "Second page text" not in result
    assert "seule servie" in result
    assert "selector='2-2'" in result


def test_pdf_read_with_range_selector(tmp_path):
    path = _make_pdf(tmp_path, ["Page A", "Page B", "Page C"])
    result = read_document("pdf", path, "2-3")
    assert "Page B" in result
    assert "Page C" in result
    assert "Page A" not in result


def test_pdf_read_scanned_page_reports_no_ocr_note(tmp_path):
    path = _make_pdf(tmp_path, [None])  # page sans texte, simulateur de scan
    result = read_document("pdf", path, "1")
    assert "OCR" in result


def test_pdf_search_finds_match_and_labels_by_page(tmp_path):
    from mcp_docs.formats import pdf_search
    from mcp_docs.search import parse_query

    path = _make_pdf(tmp_path, ["Le chat noir dort", "Rien ici"])
    results = pdf_search(path, parse_query("chat noir"))
    assert len(results) == 1
    assert results[0].label == "1"
    assert results[0].hit_count == 2


def test_pdf_search_multi_page_hits(tmp_path):
    from mcp_docs.formats import pdf_search
    from mcp_docs.search import parse_query

    path = _make_pdf(tmp_path, ["chat noir", "chat blanc", "chien"])
    results = pdf_search(path, parse_query("chat"))
    assert [r.label for r in results] == ["1", "2"]


def test_pdf_search_no_match_returns_empty(tmp_path):
    from mcp_docs.formats import pdf_search
    from mcp_docs.search import parse_query

    path = _make_pdf(tmp_path, ["Rien de pertinent ici"])
    results = pdf_search(path, parse_query("introuvable"))
    assert results == []


def test_pdf_read_non_numeric_selector_raises_clear_error(tmp_path):
    """F1 : un selector non numérique (ex. un titre de heading docx envoyé par
    erreur sur un pdf) doit lever ToolError, pas un ValueError brut."""
    path = _make_pdf(tmp_path, ["Page one"])
    with pytest.raises(ToolError, match="selector"):
        read_document("pdf", path, "Introduction")


def test_pdf_read_empty_document_reports_clear_message(tmp_path):
    """F2 : un pdf sans page ne doit pas lever IndexError. pymupdf refuse de
    sauvegarder un document à 0 page (ValueError "cannot save with zero
    pages") — on exerce donc la garde directement via pdf_read plutôt que
    via une fixture disque, en mockant doc.page_count à 0."""
    from unittest.mock import MagicMock, patch

    from mcp_docs.formats import pdf_read

    fake_doc = MagicMock()
    fake_doc.page_count = 0
    fake_doc.__enter__ = lambda s: fake_doc
    fake_doc.__exit__ = MagicMock(return_value=False)

    with patch("fitz.open", return_value=fake_doc):
        result = pdf_read(tmp_path / "empty.pdf", None)
    assert "sans aucune page" in result


# ---------------------------------------------------------------------------
# XLSX (openpyxl)
# ---------------------------------------------------------------------------

def _make_xlsx(tmp_path, sheets):
    import openpyxl

    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(name)
        for row in rows:
            ws.append(row)
    path = tmp_path / "sheet.xlsx"
    wb.save(str(path))
    return path


def test_xlsx_list_reports_sheets_and_dimensions(tmp_path):
    path = _make_xlsx(tmp_path, {"Data": [["a", "b"], ["1", "2"]], "Empty": [["x"]]})
    result = list_document("xlsx", path)
    assert "Data" in result
    assert "Empty" in result


def test_xlsx_read_default_first_sheet(tmp_path):
    path = _make_xlsx(tmp_path, {"S1": [["h1", "h2"], ["v1", "v2"]], "S2": [["z"]]})
    result = read_document("xlsx", path, None)
    assert "h1" in result and "v1" in result
    assert "S1" in result


def test_xlsx_read_max_rows_default_truncates(tmp_path):
    rows = [[str(i)] for i in range(250)]
    path = _make_xlsx(tmp_path, {"Big": rows})
    result = read_document("xlsx", path, "Big")
    assert "Tronqué à 200 lignes" in result


def test_xlsx_read_invalid_cell_range_raises_clear_error(tmp_path):
    """F4 : une plage de cellules invalide doit lever ToolError, pas un
    ValueError openpyxl brut."""
    path = _make_xlsx(tmp_path, {"Data": [["a", "b"]]})
    with pytest.raises(ToolError, match="Plage de cellules"):
        read_document("xlsx", path, "Data!ceci n'est pas une plage")


def test_xlsx_read_unknown_sheet_raises(tmp_path):
    path = _make_xlsx(tmp_path, {"S1": [["a"]]})
    with pytest.raises(ToolError, match="Feuille inconnue"):
        read_document("xlsx", path, "Nope!A1:B2")


def test_xlsx_search_finds_match_with_cell_label(tmp_path):
    from mcp_docs.formats import xlsx_search
    from mcp_docs.search import parse_query

    path = _make_xlsx(tmp_path, {"Data": [["chat noir", "autre"], ["rien", "ici"]]})
    results = xlsx_search(path, parse_query("chat noir"))
    assert len(results) == 1
    assert results[0].label == "Data"
    assert any("Data!A1" in s for s in results[0].snippets)


def test_xlsx_search_counts_occurrences_but_one_snippet_per_cell(tmp_path):
    """D6 sur xlsx : hit_count somme les occurrences réelles (ici "chat" ×2
    dans une même cellule + ×1 dans une autre = 3), mais on garde un snippet
    par cellule matchée pour que le label reste un selector `read` atomique."""
    from mcp_docs.formats import xlsx_search
    from mcp_docs.search import parse_query

    path = _make_xlsx(tmp_path, {"Data": [["chat et chat", "chat"]]})
    results = xlsx_search(path, parse_query("chat"))
    assert len(results) == 1
    assert results[0].hit_count == 3
    assert len(results[0].snippets) == 2  # deux cellules matchées


def test_xlsx_search_round_trip_label_to_read_selector(tmp_path):
    from mcp_docs.formats import xlsx_search
    from mcp_docs.search import parse_query

    path = _make_xlsx(tmp_path, {"Data": [["x", "chat noir"]]})
    results = xlsx_search(path, parse_query("chat noir"))
    snippet = results[0].snippets[0]
    coord = snippet.split(" : ")[0]  # ex. "Data!B1"
    result = read_document("xlsx", path, coord)
    assert "chat noir" in result


def test_xlsx_search_multi_sheet_hits(tmp_path):
    from mcp_docs.formats import xlsx_search
    from mcp_docs.search import parse_query

    path = _make_xlsx(tmp_path, {"S1": [["chat"]], "S2": [["chat"]], "S3": [["chien"]]})
    results = xlsx_search(path, parse_query("chat"))
    assert {r.label for r in results} == {"S1", "S2"}


def test_xlsx_search_beyond_read_row_cap_still_found(tmp_path):
    """Régression : search ne doit PAS hériter de MAX_XLSX_ROWS_DEFAULT (200) —
    un match à la ligne 300 doit rester trouvable, contrairement à read sans
    selector explicite sur toute la plage."""
    from mcp_docs.formats import xlsx_search
    from mcp_docs.search import parse_query

    rows = [["rien"] for _ in range(299)] + [["chat noir"]]
    path = _make_xlsx(tmp_path, {"Big": rows})
    results = xlsx_search(path, parse_query("chat noir"))
    assert len(results) == 1
    assert any("Big!A300" in s for s in results[0].snippets)


def test_xlsx_search_does_not_match_formula_text_with_data_only(tmp_path):
    """data_only=True lit la valeur calculée en cache, pas la formule — une
    requête sur le texte de la formule elle-même ne peut pas matcher (même
    limitation documentée que xlsx_read)."""
    import openpyxl

    from mcp_docs.formats import xlsx_search
    from mcp_docs.search import parse_query

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Calc"
    ws["A1"] = "=SUM(1,2)"
    path = tmp_path / "formula.xlsx"
    wb.save(str(path))

    results = xlsx_search(path, parse_query("SUM"))
    assert results == []


# ---------------------------------------------------------------------------
# DOCX (python-docx)
# ---------------------------------------------------------------------------

def _make_docx(tmp_path, structure, tables=None):
    """structure: liste de (level_or_None, text) ; level=None -> paragraphe normal.
    tables (optionnel) : liste de tables, chaque table = liste de lignes,
    chaque ligne = liste de textes de cellule."""
    import docx

    d = docx.Document()
    for level, text in structure:
        if level is None:
            d.add_paragraph(text)
        else:
            d.add_heading(text, level=level)
    for rows in tables or []:
        n_rows = len(rows)
        n_cols = len(rows[0]) if rows else 0
        table = d.add_table(rows=n_rows, cols=n_cols)
        for r, row_values in enumerate(rows):
            for c, value in enumerate(row_values):
                table.cell(r, c).text = value
    path = tmp_path / "doc.docx"
    d.save(str(path))
    return path


def _make_docx_with_french_heading_style(tmp_path, title_text):
    """F9 : simule un docx dont le style de heading est nommé 'Titre N'
    (nom canonique parfois exposé par un Word francophone), plutôt que le
    nom anglais 'Heading N' que add_heading() applique toujours."""
    import docx
    from docx.enum.style import WD_STYLE_TYPE

    d = docx.Document()
    styles = d.styles
    if "Titre 1" not in [s.name for s in styles]:
        french_style = styles.add_style("Titre 1", WD_STYLE_TYPE.PARAGRAPH)
        french_style.base_style = styles["Heading 1"]
    para = d.add_paragraph(title_text)
    para.style = d.styles["Titre 1"]
    path = tmp_path / "doc_fr.docx"
    d.save(str(path))
    return path


def test_docx_list_reports_french_heading_style(tmp_path):
    path = _make_docx_with_french_heading_style(tmp_path, "Introduction")
    result = list_document("docx", path)
    assert "Introduction" in result
    assert "pas de heading" not in result


def _make_docx_with_heading_style(tmp_path, style_name, title_text, filename):
    """FMT6 : même principe que _make_docx_with_french_heading_style pour une
    locale de style arbitraire (allemand, espagnol, italien)."""
    import docx
    from docx.enum.style import WD_STYLE_TYPE

    d = docx.Document()
    styles = d.styles
    if style_name not in [s.name for s in styles]:
        localized_style = styles.add_style(style_name, WD_STYLE_TYPE.PARAGRAPH)
        localized_style.base_style = styles["Heading 1"]
    para = d.add_paragraph(title_text)
    para.style = d.styles[style_name]
    path = tmp_path / filename
    d.save(str(path))
    return path


@pytest.mark.parametrize(
    "style_name,filename",
    [
        ("Überschrift 1", "doc_de.docx"),
        ("Título 1", "doc_es.docx"),
        ("Titolo 1", "doc_it.docx"),
    ],
)
def test_docx_list_reports_other_locale_heading_styles(tmp_path, style_name, filename):
    """FMT6 : allemand/espagnol/italien reconnus en plus de EN/FR."""
    path = _make_docx_with_heading_style(tmp_path, style_name, "Einleitung", filename)
    result = list_document("docx", path)
    assert "Einleitung" in result
    assert "pas de heading" not in result


def test_docx_list_reports_headings(tmp_path):
    path = _make_docx(tmp_path, [(1, "Title"), (None, "Body text"), (2, "Subtitle")])
    result = list_document("docx", path)
    assert "Title" in result
    assert "Subtitle" in result


def test_docx_list_no_headings(tmp_path):
    path = _make_docx(tmp_path, [(None, "Just a paragraph")])
    result = list_document("docx", path)
    assert "pas de heading" in result


def test_docx_read_by_heading_selector(tmp_path):
    path = _make_docx(
        tmp_path,
        [
            (1, "Intro"),
            (None, "Intro text"),
            (1, "Body"),
            (None, "Body text"),
            (1, "Conclusion"),
            (None, "Conclusion text"),
        ],
    )
    result = read_document("docx", path, "Body")
    assert "Body text" in result
    assert "Intro text" not in result
    assert "Conclusion text" not in result


def test_docx_read_unknown_heading_raises(tmp_path):
    path = _make_docx(tmp_path, [(1, "Intro"), (None, "Text")])
    with pytest.raises(ToolError, match="introuvable"):
        read_document("docx", path, "Nonexistent")


def test_docx_list_reports_tables(tmp_path):
    path = _make_docx(tmp_path, [], tables=[[["a", "b"], ["c", "d"]]])
    result = list_document("docx", path)
    assert "Tableaux (1)" in result
    assert "2 ligne(s) x 2 colonne(s)" in result


def test_docx_read_tabular_document_not_empty(tmp_path):
    """Document sans aucun paragraphe non-vide, tout le contenu en table
    (formulaire questions/réponses) : `read` sans selector doit rendre le
    contenu, pas une chaîne vide."""
    path = _make_docx(
        tmp_path, [],
        tables=[[["1", "Combien de branches ?"], ["2", "Pourquoi le logo est bleu ?"]]],
    )
    result = read_document("docx", path, None)
    assert "Combien de branches ?" in result
    assert "Pourquoi le logo est bleu ?" in result


def test_docx_read_paragraphs_and_tables_both_present(tmp_path):
    path = _make_docx(
        tmp_path, [(1, "Title"), (None, "Intro paragraph")],
        tables=[[["k", "v"]]],
    )
    result = read_document("docx", path, None)
    assert "Intro paragraph" in result
    assert "k | v" in result


def test_docx_read_by_heading_selector_ignores_tables(tmp_path):
    """Le mode selector-heading reste borné aux paragraphes : une table présente
    ailleurs dans le document n'est pas mélangée à une section ciblée."""
    path = _make_docx(
        tmp_path,
        [(1, "Body"), (None, "Body text")],
        tables=[[["unrelated", "table"]]],
    )
    result = read_document("docx", path, "Body")
    assert "Body text" in result
    assert "unrelated" not in result


def test_docx_search_finds_match_in_named_section(tmp_path):
    from mcp_docs.formats import docx_search
    from mcp_docs.search import parse_query

    path = _make_docx(
        tmp_path,
        [
            (1, "Intro"),
            (None, "Rien ici"),
            (1, "Chats"),
            (None, "Le chat noir dort"),
        ],
    )
    results = docx_search(path, parse_query("chat noir"))
    assert len(results) == 1
    assert results[0].label == "Chats"


def test_docx_search_round_trip_title_to_selector(tmp_path):
    from mcp_docs.formats import docx_search
    from mcp_docs.search import parse_query

    path = _make_docx(tmp_path, [(1, "Chats"), (None, "Le chat noir dort")])
    results = docx_search(path, parse_query("chat noir"))
    label = results[0].label
    result = read_document("docx", path, label)
    assert "chat noir" in result.lower()


def test_docx_search_matches_preamble_before_first_heading(tmp_path):
    from mcp_docs.formats import docx_search
    from mcp_docs.search import parse_query

    path = _make_docx(tmp_path, [(None, "Chat noir avant tout heading"), (1, "Suite"), (None, "Rien")])
    results = docx_search(path, parse_query("chat noir"))
    assert len(results) == 1
    assert results[0].label == "(préambule)"


def test_docx_search_document_without_heading_uses_corps_label(tmp_path):
    from mcp_docs.formats import docx_search
    from mcp_docs.search import parse_query

    path = _make_docx(tmp_path, [(None, "Chat noir sans structure de titres")])
    results = docx_search(path, parse_query("chat noir"))
    assert len(results) == 1
    assert results[0].label == "(corps)"


def test_docx_search_matches_table_content(tmp_path):
    from mcp_docs.formats import docx_search
    from mcp_docs.search import parse_query

    path = _make_docx(tmp_path, [(1, "Titre")], tables=[[["chat noir", "autre"]]])
    results = docx_search(path, parse_query("chat noir"))
    assert any(r.label == "(tableaux)" for r in results)


def test_docx_search_multiple_terms_in_same_section(tmp_path):
    """hit_count compte toutes les occurrences des termes dans l'unité (un
    snippet par occurrence), pas le nombre de termes distincts matchés."""
    from mcp_docs.formats import docx_search
    from mcp_docs.search import parse_query

    path = _make_docx(tmp_path, [(1, "Chats"), (None, "chat noir"), (None, "chat blanc")])
    results = docx_search(path, parse_query("noir blanc"))
    assert len(results) == 1
    # "noir" ×1 + "blanc" ×1 = 2 occurrences (une section unique concaténant tout).
    assert results[0].hit_count == 2


def test_docx_search_counts_repeated_occurrences(tmp_path):
    """Un terme répété plusieurs fois dans la même unité compte pour autant
    d'occurrences (régression D6 : avant, seule la première comptait)."""
    from mcp_docs.formats import docx_search
    from mcp_docs.search import parse_query

    path = _make_docx(
        tmp_path,
        [(1, "Animaux"), (None, "chat noir"), (None, "chat blanc"), (None, "chat gris")],
    )
    results = docx_search(path, parse_query("chat"))
    assert len(results) == 1
    assert results[0].hit_count == 3
    assert len(results[0].snippets) == 3


def test_docx_sections_factorization_does_not_break_docx_read(tmp_path):
    """La factorisation potentielle de docx_search via _docx_sections ne doit
    pas régresser docx_read : contrôle croisé sur un document mixte."""
    path = _make_docx(
        tmp_path,
        [(1, "Intro"), (None, "Intro text"), (1, "Body"), (None, "Body text")],
    )
    assert "Body text" in read_document("docx", path, "Body")
    assert "Intro text" not in read_document("docx", path, "Body")


# ---------------------------------------------------------------------------
# PPTX (python-pptx)
# ---------------------------------------------------------------------------

def _make_pptx(tmp_path, slide_titles):
    from pptx import Presentation

    prs = Presentation()
    for title in slide_titles:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
    path = tmp_path / "deck.pptx"
    prs.save(str(path))
    return path


def test_pptx_list_reports_slide_count_and_titles(tmp_path):
    path = _make_pptx(tmp_path, ["Intro", "Details", "Thanks"])
    result = list_document("pptx", path)
    assert "3 slide" in result
    assert "Intro" in result
    assert "Details" in result
    assert "Thanks" in result


def test_pptx_read_without_selector_first_slide_only(tmp_path):
    """F6 : idem pdf, notice de pagination par slide plutôt que "Tronqué"."""
    path = _make_pptx(tmp_path, ["Intro", "Details"])
    result = read_document("pptx", path, None)
    assert "Intro" in result
    assert "seule servie" in result
    assert "selector='2-2'" in result


def test_pptx_read_with_range(tmp_path):
    path = _make_pptx(tmp_path, ["One", "Two", "Three"])
    result = read_document("pptx", path, "2-3")
    assert "Two" in result
    assert "Three" in result
    assert "One" not in result


def test_pptx_search_finds_match_and_labels_by_slide(tmp_path):
    from mcp_docs.formats import pptx_search
    from mcp_docs.search import parse_query

    path = _make_pptx(tmp_path, ["Chat noir", "Autre slide"])
    results = pptx_search(path, parse_query("chat"))
    assert len(results) == 1
    assert results[0].label == "1"


def test_pptx_search_multi_slide_hits(tmp_path):
    from mcp_docs.formats import pptx_search
    from mcp_docs.search import parse_query

    path = _make_pptx(tmp_path, ["Chat noir", "Chat blanc", "Chien"])
    results = pptx_search(path, parse_query("chat"))
    assert [r.label for r in results] == ["1", "2"]


def test_pptx_search_no_match_returns_empty(tmp_path):
    from mcp_docs.formats import pptx_search
    from mcp_docs.search import parse_query

    path = _make_pptx(tmp_path, ["Rien de pertinent"])
    results = pptx_search(path, parse_query("introuvable"))
    assert results == []


def test_pptx_read_non_numeric_selector_raises_clear_error(tmp_path):
    """F1 : même garde que pdf_read pour un selector non numérique."""
    path = _make_pptx(tmp_path, ["Intro"])
    with pytest.raises(ToolError, match="selector"):
        read_document("pptx", path, "Introduction")


def test_pptx_read_empty_document_reports_clear_message(tmp_path):
    """F2 : une pptx sans slide (le template par défaut en a zéro) ne doit
    pas lever IndexError."""
    from pptx import Presentation

    prs = Presentation()
    path = tmp_path / "empty.pptx"
    prs.save(str(path))
    result = read_document("pptx", path, None)
    assert "sans aucune slide" in result


# ---------------------------------------------------------------------------
# ZIP (stdlib zipfile)
# ---------------------------------------------------------------------------

def _make_zip(tmp_path, members):
    import zipfile

    path = tmp_path / "archive.zip"
    with zipfile.ZipFile(path, "w") as z:
        for name, content in members.items():
            z.writestr(name, content)
    return path


def test_zip_list_reports_entries(tmp_path):
    path = _make_zip(tmp_path, {"a.txt": "hello", "dir/b.txt": "world"})
    result = list_document("zip", path)
    assert "a.txt" in result
    assert "dir/b.txt" in result


def test_zip_read_member_returns_text_content(tmp_path):
    path = _make_zip(tmp_path, {"a.txt": "hello world"})
    result = zip_read_member(path, "a.txt")
    assert "hello world" in result


def test_zip_read_member_unknown_raises(tmp_path):
    path = _make_zip(tmp_path, {"a.txt": "hello"})
    with pytest.raises(ToolError, match="introuvable"):
        zip_read_member(path, "missing.txt")


def test_zip_read_binary_member_reports_note_instead_of_garbage(tmp_path):
    path = _make_zip(tmp_path, {})
    import zipfile

    with zipfile.ZipFile(path, "a") as z:
        z.writestr("blob.bin", b"\xff\xfe\x00\x01binary")
    result = zip_read_member(path, "blob.bin")
    assert "binaire" in result


def test_read_document_zip_without_path_raises():
    with pytest.raises(ToolError, match="path"):
        read_document("zip", Path("/nonexistent"), None)


def _make_corrupt_zip(tmp_path):
    """Magic bytes PK valides (détecté comme 'zip' par sniff), mais contenu
    tronqué/corrompu — zipfile.ZipFile() doit lever BadZipFile à l'ouverture."""
    path = tmp_path / "corrupt.zip"
    path.write_bytes(b"PK\x03\x04" + b"\x00" * 20)
    return path


def test_zip_list_corrupt_archive_raises_clear_error(tmp_path):
    """F3 : BadZipFile ne doit jamais fuiter brute — docs/servers.md promet une
    erreur claire pour toute incohérence de zip (R5)."""
    path = _make_corrupt_zip(tmp_path)
    with pytest.raises(ToolError, match="corrompue"):
        list_document("zip", path)


def test_zip_search_corrupt_archive_raises_clear_error(tmp_path):
    from mcp_docs.formats import zip_search
    from mcp_docs.search import parse_query

    path = _make_corrupt_zip(tmp_path)
    with pytest.raises(ToolError, match="corrompue"):
        zip_search(path, parse_query("chat"))


# ---------------------------------------------------------------------------
# zip_search — membres texte uniquement
# ---------------------------------------------------------------------------

def test_zip_search_finds_match_in_text_member(tmp_path):
    from mcp_docs.formats import zip_search
    from mcp_docs.search import parse_query

    path = _make_zip(tmp_path, {"a.txt": "Le chat noir dort", "b.txt": "rien ici"})
    results = zip_search(path, parse_query("chat noir"))
    labeled = [r for r in results if r.label == "a.txt"]
    assert len(labeled) == 1


def test_zip_search_round_trip_label_to_read_path(tmp_path):
    from mcp_docs.formats import zip_search
    from mcp_docs.search import parse_query

    path = _make_zip(tmp_path, {"a.txt": "Le chat noir dort"})
    results = zip_search(path, parse_query("chat noir"))
    label = next(r.label for r in results if r.label != "(note)")
    result = zip_read_member(path, label)
    assert "chat noir" in result.lower()


def test_zip_search_targeted_member_via_path(tmp_path):
    from mcp_docs.formats import zip_search
    from mcp_docs.search import parse_query

    path = _make_zip(tmp_path, {"a.txt": "chat", "b.txt": "chat"})
    results = zip_search(path, parse_query("chat"), member_path="a.txt")
    assert [r.label for r in results] == ["a.txt"]


def test_zip_search_targeted_binary_member_raises_clear_error(tmp_path):
    """F5/FMT7 : cibler explicitement un membre binaire via member_path doit
    lever une erreur claire (réponse définitive et actionnable), pas un
    UnitResult à 0 hit qui se lirait comme une recherche effectuée sans
    résultat ("0 occurrence(s) dans 1 unité(s)")."""
    from mcp_docs.formats import zip_search
    from mcp_docs.search import parse_query
    from mcp_docs.session import ToolError

    path = _make_zip(tmp_path, {})
    import zipfile

    with zipfile.ZipFile(path, "a") as z:
        z.writestr("blob.bin", b"\xff\xfe\x00\x01binary")

    with pytest.raises(ToolError, match="non cherché"):
        zip_search(path, parse_query("chat"), member_path="blob.bin")


def test_zip_search_targeted_structured_member_raises_clear_error(tmp_path):
    """F5/FMT7 : idem pour un membre reconnu comme document structuré (docx)."""
    from mcp_docs.formats import zip_search
    from mcp_docs.search import parse_query
    from mcp_docs.session import ToolError

    inner_docx = _make_docx(tmp_path, [(None, "chat noir")])
    path = _make_zip(tmp_path, {})
    import zipfile

    with zipfile.ZipFile(path, "a") as z:
        z.write(inner_docx, "inner.docx")

    with pytest.raises(ToolError, match="non cherché"):
        zip_search(path, parse_query("chat"), member_path="inner.docx")


def test_zip_search_structured_and_binary_members_ignored_and_noted(tmp_path):
    from mcp_docs.formats import zip_search
    from mcp_docs.search import parse_query

    path = _make_zip(tmp_path, {"a.txt": "chat noir"})
    import zipfile

    with zipfile.ZipFile(path, "a") as z:
        z.writestr("blob.bin", b"\xff\xfe\x00\x01binary")

    inner_docx = _make_docx(tmp_path, [(None, "chat noir dans un docx imbriqué")])
    with zipfile.ZipFile(path, "a") as z:
        z.writestr("nested.docx", inner_docx.read_bytes())

    results = zip_search(path, parse_query("chat noir"))
    note = next(r for r in results if r.label == "(note)")
    assert "blob.bin" in note.snippets[0]
    assert "nested.docx" in note.snippets[0]
    # Le membre docx imbriqué n'est pas cherché malgré son contenu matchable.
    assert not any(r.label == "nested.docx" for r in results)


def test_zip_search_cumulative_size_budget_stops_scan(tmp_path, monkeypatch):
    """FMT5 : le garde par membre (MAX_UNZIP_MB) ne bornait pas le cumul —
    plusieurs membres, chacun sous la limite individuelle, pouvaient être
    balayés en entier sans borne globale. Un budget cumulé abaisse le seuil
    pour le test et vérifie que le balayage s'arrête avant le dernier membre."""
    import mcp_docs.formats as fmt

    monkeypatch.setattr(fmt, "MAX_UNZIP_MB", 100)
    monkeypatch.setattr(fmt, "_ZIP_SEARCH_TIME_BUDGET_S", 45)
    # Trois membres de ~0.6 Mo chacun (sous MAX_UNZIP_MB par membre) ; on
    # abaisse le budget cumulé à 1 Mo pour forcer l'arrêt avant le 3e.
    content = "chat noir " * 60000  # ~0.6 Mo
    path = _make_zip(tmp_path, {"a.txt": content, "b.txt": content, "c.txt": content})

    monkeypatch.setattr(fmt, "MAX_UNZIP_MB", 1)
    results = fmt.zip_search(path, parse_query("chat"))
    note = next((r for r in results if r.label == "(note)"), None)
    assert note is not None
    # Budget de temps laissé à 45 s : c'est bien le cumul qui coupe, la note
    # doit le nommer précisément (sinon le test passerait pour la mauvaise raison).
    assert "budget cumulé" in note.snippets[0]
    assert "c.txt" in note.snippets[0]


def test_zip_search_time_budget_stops_scan(tmp_path, monkeypatch):
    """FMT5 : budget de temps global (MIAOU_DOCS_SCAN_TIMEOUT_S). Budget mis à
    0 s → le balayage s'arrête au premier membre, avec une note nommant le
    budget de temps et le membre non couvert, sans lever d'erreur."""
    import mcp_docs.formats as fmt

    monkeypatch.setattr(fmt, "_ZIP_SEARCH_TIME_BUDGET_S", 0)
    path = _make_zip(tmp_path, {"a.txt": "chat noir", "b.txt": "chat noir"})

    results = fmt.zip_search(path, parse_query("chat"))
    note = next((r for r in results if r.label == "(note)"), None)
    assert note is not None
    assert "budget de temps" in note.snippets[0]
    assert "a.txt" in note.snippets[0]


def test_zip_search_time_budget_is_env_overridable(monkeypatch):
    """FMT5 : le budget de temps suit le même mécanisme _env_int que les autres
    caps de mcp_docs (défaut 30 s, surchargeable par MIAOU_DOCS_SCAN_TIMEOUT_S)."""
    import mcp_docs.session as docs_session

    monkeypatch.setenv("MIAOU_DOCS_SCAN_TIMEOUT_S", "5")
    assert docs_session._env_int("MIAOU_DOCS_SCAN_TIMEOUT_S", 30) == 5
    monkeypatch.delenv("MIAOU_DOCS_SCAN_TIMEOUT_S")
    assert docs_session._env_int("MIAOU_DOCS_SCAN_TIMEOUT_S", 30) == 30


def test_zip_search_nominal_scan_unaffected_by_budgets(tmp_path):
    """FMT5 non-régression : une archive ordinaire reste balayée en entier, tous
    les membres cherchés, sans note de zone aveugle (budgets par défaut)."""
    import mcp_docs.formats as fmt

    path = _make_zip(
        tmp_path,
        {"a.txt": "chat noir", "b.txt": "chat gris", "c.txt": "rien ici"},
    )
    results = fmt.zip_search(path, parse_query("chat"))
    labels = {r.label for r in results}
    assert labels == {"a.txt", "b.txt"}
    assert "(note)" not in labels


def test_zip_search_no_match_returns_empty_without_note(tmp_path):
    from mcp_docs.formats import zip_search
    from mcp_docs.search import parse_query

    path = _make_zip(tmp_path, {"a.txt": "rien de pertinent"})
    results = zip_search(path, parse_query("introuvable"))
    assert results == []


def test_zip_search_rejects_traversal_member_when_targeted(tmp_path):
    from mcp_docs.formats import zip_search
    from mcp_docs.search import parse_query

    path = _make_raw_zip(tmp_path, {"../evil.txt": "chat noir"})
    with pytest.raises(ToolError, match="traversal"):
        zip_search(path, parse_query("chat"), member_path="../evil.txt")


def test_zip_search_skips_traversal_member_in_full_scan(tmp_path):
    from mcp_docs.formats import zip_search
    from mcp_docs.search import parse_query

    path = _make_raw_zip(tmp_path, {"../evil.txt": "chat noir", "safe.txt": "chat noir"})
    results = zip_search(path, parse_query("chat noir"))
    assert any(r.label == "safe.txt" for r in results)
    assert not any(r.label == "../evil.txt" for r in results)


# ---------------------------------------------------------------------------
# docx/xlsx/pdf/pptx détectés par sniff interne d'un zip (Office = zip signé)
# ---------------------------------------------------------------------------

def test_detect_kind_docx_sniffed_from_zip_signature(tmp_path):
    path = _make_docx(tmp_path, [(None, "hi")])
    assert detect_kind(path) == "docx"


def test_detect_kind_xlsx_sniffed_from_zip_signature(tmp_path):
    path = _make_xlsx(tmp_path, {"S1": [["a"]]})
    assert detect_kind(path) == "xlsx"


def test_detect_kind_pptx_sniffed_from_zip_signature(tmp_path):
    path = _make_pptx(tmp_path, ["Slide"])
    assert detect_kind(path) == "pptx"


# ---------------------------------------------------------------------------
# Bout en bout via les outils MCP (list/read réels, pas seulement formats.py)
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_list_tool_end_to_end_pdf(workdir):
    from mcp_docs import server as docs_server

    path = _make_pdf(workdir, ["irrelevant, materialisation manuelle ci-dessous"])
    # Matérialise directement dans le cache de session pour éviter de repasser par b64 ici.
    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(path.read_bytes())

    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool("list", {"ref": "att-1", "session_id": "conv-1"})
    assert "PDF" in result


@pytest.mark.asyncio
async def test_list_tool_end_to_end_resource_ref(workdir):
    """res_<id> (ressource de session côté client, MIAOU lot K) doit être
    accepté et matérialisable/lisible de bout en bout comme att-N/file-<id>."""
    from mcp_docs import server as docs_server

    path = _make_pdf(workdir, ["irrelevant, materialisation manuelle ci-dessous"])
    dest = session_dir("conv-1") / "res_abc123.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(path.read_bytes())

    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool("list", {"ref": "res_abc123", "session_id": "conv-1"})
    assert "PDF" in result


def test_list_tool_docstring_does_not_restrict_ref_to_att_n():
    """DOC4 : la docstring de `list` ne doit pas décrire `ref` comme un
    att-N — trois familles existent (att-N/file-<id>/res_<id>), un modèle
    scrupuleux pourrait sinon refuser de passer un handle res_… à `list`."""
    from mcp_docs import server as docs_server

    tm = docs_server.mcp._tool_manager
    tool = next(t for t in tm.list_tools() if t.name == "list")
    assert "att-N" not in tool.description


@pytest.mark.asyncio
async def test_read_tool_end_to_end_resource_ref(workdir):
    from mcp_docs import server as docs_server

    path = _make_pdf(workdir, ["Resource ref page text"])
    dest = session_dir("conv-1") / "res_abc123.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(path.read_bytes())

    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool("read", {"ref": "res_abc123", "session_id": "conv-1"})
    assert "Resource ref page text" in result


@pytest.mark.asyncio
async def test_read_tool_requires_path_for_zip(workdir):
    from mcp.server.fastmcp.exceptions import ToolError as FastMCPToolError

    from mcp_docs import server as docs_server

    path = _make_zip(workdir, {"a.txt": "hi"})
    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(path.read_bytes())

    tm = docs_server.mcp._tool_manager
    with pytest.raises(FastMCPToolError, match="path"):
        await tm.call_tool("read", {"ref": "att-1", "session_id": "conv-1"})


# ---------------------------------------------------------------------------
# D3 — Sécurité archives : zip-slip, tailles, chiffrement, imbrication
# ---------------------------------------------------------------------------

def _make_raw_zip(tmp_path, entries, name="raw.zip"):
    """entries: dict[str, bytes|str]. writestr accepte str ou bytes."""
    import zipfile

    path = tmp_path / name
    with zipfile.ZipFile(path, "w") as z:
        for member_name, content in entries.items():
            z.writestr(member_name, content)
    return path


def test_zip_read_member_rejects_relative_traversal(tmp_path):
    path = _make_raw_zip(tmp_path, {"../../etc/passwd": "pwned"})
    with pytest.raises(ToolError, match="traversal"):
        zip_read_member(path, "../../etc/passwd")


def test_zip_read_member_rejects_absolute_path(tmp_path):
    path = _make_raw_zip(tmp_path, {"/absolute/path.txt": "pwned"})
    with pytest.raises(ToolError, match="traversal"):
        zip_read_member(path, "/absolute/path.txt")


def test_zip_list_flags_traversal_entries_without_blocking_list(tmp_path):
    """list reste possible sur une archive contenant des entrées suspectes —
    seule l'extraction (read) est bloquée."""
    path = _make_raw_zip(tmp_path, {"../evil.txt": "x", "safe.txt": "ok"})
    result = list_document("zip", path)
    assert "chemin suspect" in result
    assert "safe.txt" in result


def _make_encrypted_zip(tmp_path):
    """stdlib zipfile ne sait pas écrire de zip chiffré — dépend du binaire 'zip'."""
    import subprocess

    src = tmp_path / "secret.txt"
    src.write_text("top secret")
    archive = tmp_path / "encrypted.zip"
    subprocess.run(
        ["zip", "-P", "testpass", "-j", str(archive), str(src)],
        check=True,
        capture_output=True,
    )
    return archive


@_requires_zip_cli
def test_zip_read_member_rejects_encrypted(tmp_path):
    import zipfile

    archive = _make_encrypted_zip(tmp_path)
    with zipfile.ZipFile(archive) as z:
        member_name = z.namelist()[0]

    with pytest.raises(ToolError, match="chiffré"):
        zip_read_member(archive, member_name)


@_requires_zip_cli
def test_zip_list_flags_encrypted_entries(tmp_path):
    archive = _make_encrypted_zip(tmp_path)
    result = list_document("zip", archive)
    assert "chiffré" in result


def test_zip_read_member_rejects_oversized_declared_header(tmp_path):
    """La garde sur ZipInfo.file_size (en-tête) rejette avant même d'ouvrir le flux."""
    path = _make_raw_zip(tmp_path, {"big.txt": "x" * 1000})
    with patch("mcp_docs.formats.MAX_UNZIP_MB", 0):
        with pytest.raises(ToolError, match="volumineux"):
            zip_read_member(path, "big.txt")


def test_zip_read_member_normal_stream_within_limit_succeeds(tmp_path):
    """Garde-fou : un membre de taille normale, sous la limite, reste lisible
    (la garde en flux ne doit pas produire de faux positif)."""
    path = _make_raw_zip(tmp_path, {"ok.txt": "A" * 1000})
    result = zip_read_member(path, "ok.txt")
    assert "A" * 100 in result


def test_zip_read_member_nested_archive_reported_not_extracted(tmp_path):
    import zipfile

    inner = tmp_path / "inner.zip"
    with zipfile.ZipFile(inner, "w") as z:
        z.writestr("a.txt", "hi")

    outer = _make_raw_zip(tmp_path, {})
    with zipfile.ZipFile(outer, "a") as z:
        z.write(inner, arcname="nested.zip")

    result = zip_read_member(outer, "nested.zip")
    assert "imbriquée" in result


def test_zip_list_flags_nested_archive_candidates(tmp_path):
    import zipfile

    inner = tmp_path / "inner.zip"
    with zipfile.ZipFile(inner, "w") as z:
        z.writestr("a.txt", "hi")

    outer = _make_raw_zip(tmp_path, {})
    with zipfile.ZipFile(outer, "a") as z:
        z.write(inner, arcname="nested.zip")

    result = list_document("zip", outer)
    assert "archive imbriquée potentielle" in result


def test_zip_list_reports_total_declared_size_over_limit(tmp_path):
    path = _make_raw_zip(tmp_path, {"a.txt": "x" * 1000})
    with patch("mcp_docs.formats.MAX_UNZIP_MB", 0):
        result = list_document("zip", path)
    assert "Taille décompressée totale" in result


# ---------------------------------------------------------------------------
# D-bis — extraction de documents imbriqués dans une archive (un seul niveau)
# ---------------------------------------------------------------------------

def _zip_with_member_bytes(tmp_path, member_name, member_bytes, name="outer.zip"):
    import zipfile

    path = tmp_path / name
    with zipfile.ZipFile(path, "w") as z:
        z.writestr(member_name, member_bytes)
    return path


def test_zip_read_member_docx_nested_returns_structured_content(tmp_path):
    inner = _make_docx(tmp_path, [(1, "Intro"), (None, "Nested body text")])
    outer = _zip_with_member_bytes(tmp_path, "membre.docx", inner.read_bytes())

    result = zip_read_member(outer, "membre.docx")
    assert "Nested body text" in result


def test_zip_read_member_xlsx_nested_returns_structured_content(tmp_path):
    inner = _make_xlsx(tmp_path, {"S1": [["h1", "h2"], ["v1", "v2"]]})
    outer = _zip_with_member_bytes(tmp_path, "membre.xlsx", inner.read_bytes())

    result = zip_read_member(outer, "membre.xlsx")
    assert "v1" in result


def test_zip_read_member_pptx_nested_returns_structured_content(tmp_path):
    inner = _make_pptx(tmp_path, ["Nested Slide"])
    outer = _zip_with_member_bytes(tmp_path, "membre.pptx", inner.read_bytes())

    result = zip_read_member(outer, "membre.pptx")
    assert "Nested Slide" in result


def test_zip_read_member_pdf_nested_returns_structured_content(tmp_path):
    inner = _make_pdf(tmp_path, ["Nested PDF text"])
    outer = _zip_with_member_bytes(tmp_path, "membre.pdf", inner.read_bytes())

    result = zip_read_member(outer, "membre.pdf")
    assert "Nested PDF text" in result


def test_zip_read_member_docx_nested_with_selector(tmp_path):
    inner = _make_docx(
        tmp_path,
        [(1, "Intro"), (None, "Intro text"), (1, "Body"), (None, "Body text")],
    )
    outer = _zip_with_member_bytes(tmp_path, "membre.docx", inner.read_bytes())

    result = zip_read_member(outer, "membre.docx", selector="Body")
    assert "Body text" in result
    assert "Intro text" not in result


def test_zip_list_member_docx_nested_reports_structure(tmp_path):
    inner = _make_docx(tmp_path, [(1, "Nested Title")])
    outer = _zip_with_member_bytes(tmp_path, "membre.docx", inner.read_bytes())

    result = zip_list_member(outer, "membre.docx")
    assert "Nested Title" in result


def test_zip_read_member_double_nested_zip_stays_one_level_max(tmp_path):
    """Un zip contenant un zip contenant un docx : le premier niveau (zip) doit
    rester signalé, pas de récursion au-delà d'un niveau."""
    import zipfile

    inner_docx = _make_docx(tmp_path, [(None, "should not surface")])
    inner_zip = tmp_path / "inner.zip"
    with zipfile.ZipFile(inner_zip, "w") as z:
        z.writestr("nested.docx", inner_docx.read_bytes())

    outer = _zip_with_member_bytes(tmp_path, "inner.zip", inner_zip.read_bytes())

    result = zip_read_member(outer, "inner.zip")
    assert "un seul niveau d'imbrication" in result
    assert "should not surface" not in result


def test_zip_list_member_double_nested_zip_stays_one_level_max(tmp_path):
    import zipfile

    inner_docx = _make_docx(tmp_path, [(None, "should not surface")])
    inner_zip = tmp_path / "inner.zip"
    with zipfile.ZipFile(inner_zip, "w") as z:
        z.writestr("nested.docx", inner_docx.read_bytes())

    outer = _zip_with_member_bytes(tmp_path, "inner.zip", inner_zip.read_bytes())

    result = zip_list_member(outer, "inner.zip")
    assert "un seul niveau d'imbrication" in result


def test_zip_read_member_nested_oversized_raises_clear_error(tmp_path):
    """Un membre imbriqué reconnu comme docx mais dépassant MAX_FILE_MB doit
    lever une ToolError claire, pas planter dans la lib de parsing."""
    inner = _make_docx(tmp_path, [(None, "x" * 1000)])
    outer = _zip_with_member_bytes(tmp_path, "membre.docx", inner.read_bytes())

    with patch("mcp_docs.formats.MAX_FILE_MB", 0):
        with pytest.raises(ToolError, match="volumineux"):
            zip_read_member(outer, "membre.docx")


def test_zip_read_member_nested_corrupted_docx_raises_clear_error(tmp_path):
    """Un membre dont l'extension/signature suggère un docx mais dont le
    contenu est tronqué/corrompu doit lever une ToolError, pas fuiter
    l'exception de python-docx."""
    inner = _make_docx(tmp_path, [(None, "valid docx content")])
    truncated_bytes = inner.read_bytes()[:20]  # signature zip présente, contenu invalide
    outer = _zip_with_member_bytes(tmp_path, "membre.docx", truncated_bytes)

    with pytest.raises(ToolError):
        zip_read_member(outer, "membre.docx")


# ---------------------------------------------------------------------------
# zip_extract_member_text (lot M — texte intégral, sans READ_CAP, transfert client)
# ---------------------------------------------------------------------------

def test_zip_extract_member_text_returns_full_text_no_truncation(tmp_path):
    """Un membre plus long que READ_CAP doit ressortir intégralement — aucune
    notice de troncature, contrairement à zip_read_member."""
    big = "line\n" * (docs_session.READ_CAP // 4)
    path = _make_zip(tmp_path, {"big.log": big})
    result = zip_extract_member_text(path, "big.log")
    assert result == big
    assert len(result) > docs_session.READ_CAP
    assert "Tronqué" not in result
    assert "tronqu" not in result.lower()


def test_zip_extract_member_text_rejects_traversal(tmp_path):
    path = _make_zip(tmp_path, {"a.txt": "hello"})
    with pytest.raises(ToolError, match="invalide"):
        zip_extract_member_text(path, "../../etc/passwd")


@_requires_zip_cli
def test_zip_extract_member_text_rejects_encrypted(tmp_path):
    import zipfile

    archive = _make_encrypted_zip(tmp_path)
    with zipfile.ZipFile(archive) as z:
        member_name = z.namelist()[0]

    with pytest.raises(ToolError, match="chiffré"):
        zip_extract_member_text(archive, member_name)


def test_zip_extract_member_text_rejects_structured_document(tmp_path):
    """Un membre docx doit être refusé explicitement (steer vers docs__read), pas
    renvoyé comme du texte (XML brut de l'archive Office serait incompréhensible)."""
    inner = _make_docx(tmp_path, [(1, "Intro"), (None, "Nested body text")])
    outer = _zip_with_member_bytes(tmp_path, "membre.docx", inner.read_bytes())

    with pytest.raises(ToolError, match="structuré"):
        zip_extract_member_text(outer, "membre.docx")


def test_zip_extract_member_text_rejects_binary_member(tmp_path):
    path = _make_zip(tmp_path, {})
    import zipfile

    with zipfile.ZipFile(path, "a") as z:
        z.writestr("blob.bin", b"\xff\xfe\x00\x01binary")

    with pytest.raises(ToolError, match="binaire"):
        zip_extract_member_text(path, "blob.bin")


def test_zip_extract_member_text_unknown_member_raises(tmp_path):
    path = _make_zip(tmp_path, {"a.txt": "hello"})
    with pytest.raises(ToolError, match="introuvable"):
        zip_extract_member_text(path, "missing.txt")


@pytest.mark.asyncio
async def test_extract_tool_returns_blob_resource_with_textual_mime(workdir):
    """docs__extract renvoie un TextContent descripteur (sans le texte du membre)
    + un EmbeddedResource blob mimeType text/plain — le canal K, jamais
    resource.text (le texte ne doit jamais atterrir dans le contenu textuel)."""
    from mcp_docs import server as docs_server

    outer = _make_zip(workdir, {"data.json": '{"hits": 42}'})
    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(outer.read_bytes())

    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool(
        "extract", {"ref": "att-1", "session_id": "conv-1", "path": "data.json"}
    )
    blocks = result.content if hasattr(result, "content") else result

    text_blocks = [b for b in blocks if getattr(b, "type", None) == "text"]
    resource_blocks = [b for b in blocks if getattr(b, "type", None) == "resource"]

    assert len(resource_blocks) == 1
    resource = resource_blocks[0].resource
    assert resource.mimeType == "text/plain"
    assert base64.b64decode(resource.blob).decode("utf-8") == '{"hits": 42}'

    for tb in text_blocks:
        assert "hits" not in tb.text
        assert "42" not in tb.text


@pytest.mark.asyncio
async def test_extract_tool_rejects_non_zip_ref(workdir):
    from mcp.server.fastmcp.exceptions import ToolError as FastMCPToolError

    from mcp_docs import server as docs_server

    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(b"not a zip")

    tm = docs_server.mcp._tool_manager
    with pytest.raises(FastMCPToolError, match="archive"):
        await tm.call_tool(
            "extract", {"ref": "att-1", "session_id": "conv-1", "path": "whatever.txt"}
        )


@pytest.mark.asyncio
async def test_list_tool_path_on_nested_docx_member(workdir):
    """docs__list avec path pointant vers un membre docx imbriqué renvoie sa
    structure, pas la note générique 'pas un document structuré'."""
    from mcp_docs import server as docs_server

    inner_doc = _make_docx(workdir, [(1, "Nested Title")])
    outer = _zip_with_member_bytes(workdir, "membre.docx", inner_doc.read_bytes())
    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(outer.read_bytes())

    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool(
        "list", {"ref": "att-1", "session_id": "conv-1", "path": "membre.docx"}
    )
    assert "Nested Title" in result


@pytest.mark.asyncio
async def test_read_tool_path_on_nested_docx_member(workdir):
    """docs__read avec path pointant vers un membre docx imbriqué renvoie son
    contenu structuré."""
    from mcp_docs import server as docs_server

    inner_doc = _make_docx(workdir, [(None, "Nested body via tool")])
    outer = _zip_with_member_bytes(workdir, "membre.docx", inner_doc.read_bytes())
    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(outer.read_bytes())

    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool(
        "read", {"ref": "att-1", "session_id": "conv-1", "path": "membre.docx"}
    )
    assert "Nested body via tool" in result


@pytest.mark.asyncio
async def test_list_tool_path_on_non_zip_raises(workdir):
    from mcp.server.fastmcp.exceptions import ToolError as FastMCPToolError

    from mcp_docs import server as docs_server

    path = _make_pdf(workdir, ["irrelevant"])
    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(path.read_bytes())

    tm = docs_server.mcp._tool_manager
    with pytest.raises(FastMCPToolError, match="archive"):
        await tm.call_tool("list", {"ref": "att-1", "session_id": "conv-1", "path": "x.docx"})


@pytest.mark.asyncio
async def test_search_tool_end_to_end_pdf(workdir):
    from mcp_docs import server as docs_server

    path = _make_pdf(workdir, ["Le chat noir dort", "Rien ici"])
    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(path.read_bytes())

    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool(
        "search", {"ref": "att-1", "session_id": "conv-1", "query": "chat noir"}
    )
    assert "1" in result
    assert "chat noir" in result.lower()


@pytest.mark.asyncio
async def test_search_tool_end_to_end_zip(workdir):
    from mcp_docs import server as docs_server

    path = _make_zip(workdir, {"a.txt": "Le chat noir dort", "b.txt": "rien"})
    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(path.read_bytes())

    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool(
        "search", {"ref": "att-1", "session_id": "conv-1", "query": "chat noir"}
    )
    assert "a.txt" in result
    assert "chat noir" in result.lower()


@pytest.mark.asyncio
async def test_search_ref_unknown_through_proxy_raises_jsonrpc_error(workdir):
    """search est inflatable (ref+content_b64) : REF_UNKNOWN doit aussi passer
    par le chemin JSON-RPC du proxy, pas seulement list/read."""
    import mcp.types as types
    from mcp.shared.exceptions import McpError

    from mcp_proxy import InProcessUpstream, build_proxy_server

    upstream = InProcessUpstream("mcp_docs")
    await upstream.start()

    upstreams = {"docs": upstream}
    tool_map: dict = {}
    server = build_proxy_server(upstreams, tool_map)

    list_handler = server.request_handlers[types.ListToolsRequest]
    await list_handler(types.ListToolsRequest(method="tools/list", params=None))

    call_handler = server.request_handlers[types.CallToolRequest]
    request = types.CallToolRequest(
        method="tools/call",
        params=types.CallToolRequestParams(
            name="docs__search",
            arguments={"ref": "att-99", "session_id": "conv-1", "query": "chat"},
        ),
    )

    with pytest.raises(McpError) as exc_info:
        await call_handler(request)

    assert exc_info.value.error.data["code"] == "REF_UNKNOWN"
    assert exc_info.value.error.code == REF_UNKNOWN_ERROR_CODE


# ---------------------------------------------------------------------------
# read — plage char/ligne (fenêtre glissante au-delà du cap de READ_CAP)
# ---------------------------------------------------------------------------

def test_build_range_none_when_no_params():
    from mcp_docs import _build_range

    assert _build_range(None, None, None, None) is None


def test_build_range_char_and_line_mutually_exclusive():
    from mcp_docs import _build_range

    with pytest.raises(ToolError, match="exclusifs"):
        _build_range(0, None, 1, None)


def test_build_range_char_end_without_start_rejected():
    from mcp_docs import _build_range

    with pytest.raises(ToolError, match="char_start"):
        _build_range(None, 100, None, None)


def test_build_range_line_end_without_start_rejected():
    from mcp_docs import _build_range

    with pytest.raises(ToolError, match="line_start"):
        _build_range(None, None, None, 10)


def test_build_range_line_start_must_be_one_indexed():
    from mcp_docs import _build_range

    with pytest.raises(ToolError, match="line_start"):
        _build_range(None, None, 0, None)


def test_build_range_char_start_negative_rejected():
    from mcp_docs import _build_range

    with pytest.raises(ToolError, match="char_start"):
        _build_range(-1, None, None, None)


def test_build_range_end_before_start_rejected():
    from mcp_docs import _build_range

    with pytest.raises(ToolError, match="char_end"):
        _build_range(100, 50, None, None)
    with pytest.raises(ToolError, match="line_end"):
        _build_range(None, None, 10, 5)


def test_apply_range_char_window_within_text():
    from mcp_docs.formats import TextRange, _apply_range

    body = "abcdefghij"
    text, notice = _apply_range(body, TextRange(char_start=2, char_end=5))
    assert text == "cde"


def test_apply_range_char_start_only_caps_at_read_cap(monkeypatch):
    import mcp_docs.formats as fmt

    monkeypatch.setattr(fmt, "READ_CAP", 5)
    body = "0123456789ABCDE"
    text, notice = fmt._apply_range(body, fmt.TextRange(char_start=3))
    assert text == "34567"
    assert "char_start=8" in notice  # 3 + 5 servis


def test_apply_range_char_start_beyond_end():
    from mcp_docs.formats import TextRange, _apply_range

    text, notice = _apply_range("short", TextRange(char_start=100))
    assert text == ""
    assert "au-delà" in notice


def test_apply_range_line_window():
    from mcp_docs.formats import TextRange, _apply_range

    body = "L1\nL2\nL3\nL4\nL5"
    text, notice = _apply_range(body, TextRange(line_start=2, line_end=4))
    assert text == "L2\nL3\nL4"
    assert "line_start=5" in notice  # end < total → suite


def test_apply_range_line_window_single_line_exceeds_cap_suggests_char_mode(monkeypatch):
    """F7 : quand le cap coupe avant la première fin de ligne (served_lines=0),
    suggérer line_start=même valeur n'aiderait pas à progresser — le hint doit
    pointer vers char_start."""
    import mcp_docs.formats as fmt

    monkeypatch.setattr(fmt, "READ_CAP", 5)
    body = "0123456789ABCDE"  # une seule "ligne" de 15 caractères
    text, notice = fmt._apply_range(body, fmt.TextRange(line_start=1))
    assert text == "01234"
    assert "char_start" in notice
    assert "utiliser char_start" in notice


def test_apply_range_line_start_only_to_end():
    from mcp_docs.formats import TextRange, _apply_range

    body = "L1\nL2\nL3"
    text, notice = _apply_range(body, TextRange(line_start=2))
    assert text == "L2\nL3"
    assert notice == ""  # atteint la dernière ligne


def test_apply_range_line_start_beyond_end():
    from mcp_docs.formats import TextRange, _apply_range

    body = "L1\nL2"
    text, notice = _apply_range(body, TextRange(line_start=99))
    assert text == ""
    assert "au-delà" in notice


def test_apply_range_char_start_positive_on_empty_body_reports_out_of_bounds():
    """FMT8 : char_start > 0 sur un corps vide tombait auparavant sur
    ToolError('Plage de caractères invalide : 5-None') au lieu de la notice
    hors bornes normale."""
    from mcp_docs.formats import TextRange, _apply_range

    text, notice = _apply_range("", TextRange(char_start=5))
    assert text == ""
    assert "au-delà" in notice


def test_apply_range_char_start_zero_on_empty_body_returns_empty_without_error():
    from mcp_docs.formats import TextRange, _apply_range

    text, notice = _apply_range("", TextRange(char_start=0))
    assert text == ""
    assert notice == ""


def test_parse_range_within_bounds_reports_no_notice():
    from mcp_docs.formats import _parse_range

    start, end, notice = _parse_range("2-4", 10)
    assert (start, end) == (2, 4)
    assert notice == ""


def test_parse_range_clamped_reports_notice():
    """FMT4 : une plage qui déborde du document doit désormais l'annoncer
    plutôt que de clamper silencieusement."""
    from mcp_docs.formats import _parse_range

    start, end, notice = _parse_range("5-100", 10)
    assert (start, end) == (5, 10)
    assert "5-10" in notice
    assert "5-100" in notice


def test_pdf_read_range_selector_clamped_reports_notice(tmp_path):
    """FMT4 vu depuis pdf_read : le clamp de _parse_range doit apparaître dans
    la sortie de l'outil, pas seulement dans la fonction interne."""
    path = _make_pdf(tmp_path, ["Page A", "Page B", "Page C"])
    result = read_document("pdf", path, "2-100")
    assert "Plage ramenée à 2-3" in result


def test_pdf_read_char_range_exceeds_cap_via_offset(tmp_path, monkeypatch):
    """Une page unique dépassant READ_CAP devient lisible en entier via
    plusieurs appels char_start décalés. La plage porte sur le body rendu
    complet, en-tête '--- Page N ---' compris."""
    import mcp_docs.formats as fmt

    monkeypatch.setattr(fmt, "READ_CAP", 40)
    big = "X" * 200
    path = _make_pdf(tmp_path, [big])

    # En-tête "--- Page 1 ---\n" = 15 caractères, puis les X. Un cap à 40 sert
    # 40 caractères de fenêtre (en-tête + 25 X) et annonce l'offset suivant.
    first = fmt.read_document("pdf", path, "1", fmt.TextRange(char_start=0))
    assert first.count("X") == 25
    assert "char_start=40" in first

    second = fmt.read_document("pdf", path, "1", fmt.TextRange(char_start=40))
    assert second.count("X") == 40
    assert "char_start=80" in second


def test_read_document_xlsx_rejects_range(tmp_path):
    from mcp_docs.formats import TextRange, read_document

    path = _make_xlsx(tmp_path, {"S1": [["a", "b"]]})
    with pytest.raises(ToolError, match="xlsx"):
        read_document("xlsx", path, None, TextRange(char_start=0))


def test_zip_read_member_text_char_range(tmp_path):
    from mcp_docs.formats import TextRange, zip_read_member

    path = _make_zip(tmp_path, {"a.txt": "abcdefghij"})
    result = zip_read_member(path, "a.txt", rng=TextRange(char_start=2, char_end=5))
    assert result.startswith("cde")


@pytest.mark.asyncio
async def test_read_tool_char_range_end_to_end(workdir, monkeypatch):
    import mcp_docs.formats as fmt

    monkeypatch.setattr(fmt, "READ_CAP", 30)
    from mcp_docs import server as docs_server

    path = _make_pdf(workdir, ["Y" * 100])
    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(path.read_bytes())

    tm = docs_server.mcp._tool_manager
    result = await tm.call_tool(
        "read",
        {"ref": "att-1", "session_id": "conv-1", "selector": "1", "char_start": 50},
    )
    assert "Y" in result
    assert "char_start" in result


@pytest.mark.asyncio
async def test_read_tool_rejects_mixed_range_modes(workdir):
    from mcp.server.fastmcp.exceptions import ToolError as FastMCPToolError

    from mcp_docs import server as docs_server

    path = _make_pdf(workdir, ["text"])
    dest = session_dir("conv-1") / "att-1.bin"
    dest.parent.mkdir(parents=True, exist_ok=True)
    dest.write_bytes(path.read_bytes())

    tm = docs_server.mcp._tool_manager
    with pytest.raises(FastMCPToolError, match="exclusif"):
        await tm.call_tool(
            "read",
            {"ref": "att-1", "session_id": "conv-1", "char_start": 0, "line_start": 1},
        )


# ---------------------------------------------------------------------------
# search.py — logique pure (fold, parse_query, match_unit, make_snippet)
# Aucune fixture binaire ici : ces tests n'importent aucune lib de doc.
# ---------------------------------------------------------------------------

def test_fold_lowercases_and_strips_accents():
    assert fold("Été") == "ete"
    assert fold("ÀÉÈÊÇÙÎÔ") == "aeeecuio"


def test_fold_ligatures_oe_ae():
    assert fold("sœur") == "soeur"
    assert fold("nœud") == "noeud"
    assert fold("Œuvre") == "oeuvre"
    assert fold("Æquo") == "aequo"


def test_fold_case_insensitive_equivalence():
    assert fold("Café") == fold("CAFÉ") == fold("cafe")


def test_parse_query_rejects_empty():
    with pytest.raises(ToolError):
        parse_query("")
    with pytest.raises(ToolError):
        parse_query("   ")


def test_parse_query_whitespace_terms():
    q = parse_query("chat noir")
    assert q.terms == ["chat", "noir"]


def test_parse_query_quoted_phrase():
    q = parse_query('"chat noir" hiver')
    assert q.terms == ["chat noir", "hiver"]


def test_parse_query_unclosed_quote_falls_back_to_terms():
    q = parse_query('foo "bar baz')
    assert q.terms == ["foo", "bar", "baz"]


def test_match_unit_all_terms_present():
    q = parse_query("chat noir")
    hits = match_unit("Le chat noir dort.", q)
    assert len(hits) == 2


def test_match_unit_partial_terms_absent_yields_no_hits():
    q = parse_query("chat souris")
    hits = match_unit("Le chat noir dort.", q)
    assert hits == []


def test_match_unit_counts_all_occurrences_of_a_term():
    q = parse_query("chat")
    hits = match_unit("Le chat, le chat, encore le chat.", q)
    assert len(hits) == 3
    # Triés par offset croissant.
    assert [h.offset for h in hits] == sorted(h.offset for h in hits)


def test_match_unit_and_gate_still_holds_with_occurrences():
    # "chat" présent 2×, "souris" absent → AND-gate invalide toute l'unité.
    q = parse_query("chat souris")
    hits = match_unit("Le chat et le chat.", q)
    assert hits == []


def test_match_unit_hits_ordered_across_terms():
    q = parse_query("chat chien")
    hits = match_unit("chien puis chat puis chien.", q)
    # 1 "chat" + 2 "chien" = 3 hits, ordonnés par position dans le texte.
    assert [h.term for h in hits] == ["chien", "chat", "chien"]


def test_match_unit_case_and_accent_insensitive():
    q = parse_query("ete")
    hits = match_unit("Il fait chaud cet Été.", q)
    assert len(hits) == 1


def test_match_unit_quoted_phrase_exact():
    q = parse_query('"chat noir"')
    hits = match_unit("Le chat noir dort, le chien blanc aussi.", q)
    assert len(hits) == 1
    hits_absent = match_unit("Le chat blanc et le chien noir.", q)
    assert hits_absent == []


def test_make_snippet_middle_of_text():
    text = "a" * 200 + "MATCH" + "b" * 200
    snippet = make_snippet(text, offset=200, length=5, radius=10)
    assert snippet.startswith("…")
    assert snippet.endswith("…")
    assert "MATCH" in snippet


def test_make_snippet_at_start_no_leading_ellipsis():
    text = "MATCH" + "b" * 200
    snippet = make_snippet(text, offset=0, length=5, radius=10)
    assert not snippet.startswith("…")
    assert snippet.endswith("…")


def test_make_snippet_at_end_no_trailing_ellipsis():
    text = "a" * 200 + "MATCH"
    snippet = make_snippet(text, offset=200, length=5, radius=10)
    assert snippet.startswith("…")
    assert not snippet.endswith("…")


def test_make_snippet_short_text_no_ellipsis():
    snippet = make_snippet("MATCH", offset=0, length=5, radius=10)
    assert snippet == "MATCH"


def test_render_results_total_reflects_real_occurrence_count():
    """Le total annoncé par render_results est la somme des hit_count (vraies
    occurrences), pas le nombre de snippets ni d'unités (D6)."""
    from mcp_docs.search import UnitResult, render_results

    units = [
        UnitResult(label="1", hit_count=3, snippets=["s1", "s2", "s3"]),
        UnitResult(label="2", hit_count=2, snippets=["s4", "s5"]),
    ]
    out = render_results("pdf", "chat", units, cap=50)
    assert "5 occurrence(s) dans 2 unité(s)" in out
    assert "--- 1 (3 occurrence(s)) ---" in out


def test_render_results_truncates_at_cap():
    from mcp_docs.search import UnitResult, render_results

    units = [UnitResult(label="1", hit_count=4, snippets=["s1", "s2", "s3", "s4"])]
    out = render_results("pdf", "chat", units, cap=2)
    assert "Tronqué à 2 snippets" in out
    assert "2 occurrence(s) supplémentaire(s)" in out


# ---------------------------------------------------------------------------
# FMT1 — la notice de troncature READ_CAP ne doit jamais être masquée par une
# notice de pagination/vide déjà posée (concaténation, pas `if not notice`).
# ---------------------------------------------------------------------------

def test_pdf_read_truncation_notice_not_masked_by_pagination_notice(tmp_path, monkeypatch):
    import mcp_docs.formats as fmt

    monkeypatch.setattr(fmt, "READ_CAP", 20)
    path = _make_pdf(tmp_path, ["X" * 200, "second page"])
    result = fmt.read_document("pdf", path, None)
    assert "seule servie" in result
    assert "Tronqué à 20 caractères" in result


def test_pptx_read_truncation_notice_not_masked_by_pagination_notice(tmp_path, monkeypatch):
    import mcp_docs.formats as fmt

    monkeypatch.setattr(fmt, "READ_CAP", 20)
    path = _make_pptx(tmp_path, ["X" * 200, "second slide"])
    result = fmt.read_document("pptx", path, None)
    assert "seule servie" in result
    assert "Tronqué à 20 caractères" in result


def test_docx_read_truncation_notice_not_masked_by_para_notice(tmp_path, monkeypatch):
    import mcp_docs.formats as fmt

    monkeypatch.setattr(fmt, "READ_CAP", 20)
    structure = [(None, f"paragraph {i}") for i in range(60)]
    path = _make_docx(tmp_path, structure)
    result = fmt.read_document("docx", path, None)
    assert "selector=<titre d'un heading>" in result
    assert "Tronqué à 20 caractères" in result


# ---------------------------------------------------------------------------
# FMT2 — exceptions des libs de parsing sur document forgé/corrompu converties
# en ToolError, sans détail interne (chemin disque, trace stdlib).
# ---------------------------------------------------------------------------

def test_pdf_read_corrupt_document_raises_clear_tool_error(tmp_path):
    path = tmp_path / "corrupt.pdf"
    path.write_bytes(b"%PDF" + b"garbage" * 20)
    with pytest.raises(ToolError, match="illisible ou corrompu") as exc_info:
        read_document("pdf", path, None)
    assert str(tmp_path) not in str(exc_info.value)


def test_docx_list_corrupt_office_zip_raises_clear_tool_error(tmp_path):
    """Zip avec un dossier word/ (détecté docx) mais sans [Content_Types].xml —
    python-docx lève une KeyError brute, doit devenir une ToolError propre."""
    import zipfile

    path = tmp_path / "broken.docx"
    with zipfile.ZipFile(path, "w") as z:
        z.writestr("word/document.xml", "<xml/>")
    with pytest.raises(ToolError, match="illisible ou corrompu"):
        list_document("docx", path)


# ---------------------------------------------------------------------------
# FMT3 — xlsx sans élément <dimension> ("unsized") : fallback plutôt que fuite
# de ValueError openpyxl.
# ---------------------------------------------------------------------------

def test_xlsx_list_unsized_sheet_falls_back(tmp_path):
    import re
    import zipfile

    src = _make_xlsx(tmp_path, {"Data": [["a", "b"], ["1", "2"]]})
    dst = tmp_path / "unsized.xlsx"
    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(dst, "w") as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "xl/worksheets/sheet1.xml":
                data = re.sub(rb"<dimension[^/]*/>", b"", data)
            zout.writestr(item, data)

    result = list_document("xlsx", dst)
    assert "Data" in result
    assert "colonne(s)" in result


# ---------------------------------------------------------------------------
# DOC1/DOC9 — matérialisation et nettoyage hors event loop (asyncio.to_thread).
# Le test vérifie à la fois que l'outil reste fonctionnel de bout en bout via
# tool_manager, et que la loop reste ordonnançable pendant l'étape bloquante
# (resolve_ref : sweep + b64decode + quota + write_bytes).
# ---------------------------------------------------------------------------

@pytest.mark.asyncio
async def test_list_materializes_off_the_event_loop(workdir, tmp_path):
    """docs__list reste fonctionnel de bout en bout, et une coroutine concurrente
    progresse pendant la matérialisation : resolve_ref (sweep + b64decode + quota
    + write_bytes) ne doit pas s'exécuter dans l'event loop."""
    import asyncio

    from mcp_docs import server as docs_server

    path = _make_zip(tmp_path, {"a.txt": "hello", "dir/b.txt": "world"})
    content_b64 = base64.b64encode(path.read_bytes()).decode()

    BLOCK = 0.3
    real_resolve = docs_session.resolve_ref

    def _slow_resolve(*args, **kwargs):
        time.sleep(BLOCK)  # gèlerait la loop si appelée hors to_thread
        return real_resolve(*args, **kwargs)

    concurrent_ran = asyncio.Event()

    async def _concurrent():
        # Se réveille bien avant la fin du resolve : n'y parvient que si la loop
        # est libre pendant l'étape bloquante.
        await asyncio.sleep(BLOCK / 3)
        concurrent_ran.set()

    async def _call():
        with patch.object(mcp_docs, "resolve_ref", _slow_resolve):
            tm = docs_server.mcp._tool_manager
            return await tm.call_tool(
                "list",
                {"ref": "att-1", "session_id": "conv-thread", "content_b64": content_b64},
            )

    call_task = asyncio.create_task(_call())
    await _concurrent()
    assert not call_task.done(), "resolve_ref trop rapide — le test ne mesure rien"
    assert concurrent_ran.is_set(), "la loop est restée bloquée pendant resolve_ref"

    result = await call_task
    assert "a.txt" in result
    assert "dir/b.txt" in result
